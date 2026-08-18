"""
ExportIntel AI — Review Presentation (guideline-aligned, 15–20 slides)
Custom professional design. Avoids generic AI-template look.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "ExportIntel_AI_Review_Presentation.pptx"

IMG_ARCH = ROOT / "ExportIntel_AI_Architecture_Diagram.png"
IMG1 = ROOT / "assets" / "image-1.png"
IMG2 = ROOT / "assets" / "image-2.png"
IMG3 = ROOT / "assets" / "image-3.png"
IMG4 = ROOT / "assets" / "image-4.png"

# Custom palette (navy + teal — not purple/cream AI defaults)
NAVY = RGBColor(15, 36, 72)
TEAL = RGBColor(14, 116, 144)
INK = RGBColor(28, 35, 45)
MUTED = RGBColor(90, 100, 112)
SOFT = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(210, 220, 230)
ACCENT = RGBColor(194, 120, 45)

SOURCE_ASSETS = [
    ROOT.parent / ".cursor" / "projects" / "c-Users-Lenovo-Desktop-Export-AI" / "assets" / "c__Users_Lenovo_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-326990d6-d86f-4625-aad4-ffff9f82faaa.png",
    ROOT.parent / ".cursor" / "projects" / "c-Users-Lenovo-Desktop-Export-AI" / "assets" / "c__Users_Lenovo_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-8ba07128-fbb0-4042-8280-29648355cab1.png",
    ROOT.parent / ".cursor" / "projects" / "c-Users-Lenovo-Desktop-Export-AI" / "assets" / "c__Users_Lenovo_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-1e0864a7-b10a-42aa-bd69-a41c45320d26.png",
    ROOT.parent / ".cursor" / "projects" / "c-Users-Lenovo-Desktop-Export-AI" / "assets" / "c__Users_Lenovo_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-d82a3a01-86ed-4cf0-bb52-7b8d5ae2353d.png",
]


def ensure_assets():
    (ROOT / "assets").mkdir(exist_ok=True)
    for src, dst in zip(SOURCE_ASSETS, [IMG1, IMG2, IMG3, IMG4]):
        if src.exists() and not dst.exists():
            dst.write_bytes(src.read_bytes())


def set_run(p, text, *, size=18, bold=False, color=INK, name="Calibri"):
    p.clear()
    r = p.add_run()
    r.text = text
    r.font.name = name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def paint_bg(slide, prs, color=SOFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # left accent rail
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), prs.slide_height)
    rail.fill.solid()
    rail.fill.fore_color.rgb = TEAL
    rail.line.fill.background()


def footer(slide, n, total, label="ExportIntel AI  ·  Product Pitch"):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(11.5), Inches(0.3))
    tf = box.text_frame
    set_run(tf.paragraphs[0], f"{label}                              {n}/{total}", size=11, color=MUTED)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT


def title_block(slide, title, subtitle=""):
    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.7))
    set_run(t.text_frame.paragraphs[0], title, size=30, bold=True, color=NAVY)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(0.95), Inches(12.2), Inches(0.4))
        set_run(s.text_frame.paragraphs[0], subtitle, size=15, color=MUTED)


def card(slide, x, y, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = LINE
    sh.adjustments[0] = 0.08
    return sh


def bullets(slide, items, x, y, w, h, *, size=17):
    card(slide, x - 0.08, y - 0.1, w + 0.16, h + 0.2)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        set_run(p, f"•  {item}", size=size, color=INK)
        p.space_after = Pt(8)
        p.level = 0


def section_hero(prs, eyebrow, title, subtitle, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs, NAVY)
    # teal strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.55), prs.slide_width, Inches(2.4))
    strip.fill.solid()
    strip.fill.fore_color.rgb = TEAL
    strip.line.fill.background()
    eb = s.shapes.add_textbox(Inches(0.8), Inches(2.75), Inches(11.5), Inches(0.4))
    set_run(eb.text_frame.paragraphs[0], eyebrow.upper(), size=14, bold=True, color=WHITE)
    tt = s.shapes.add_textbox(Inches(0.8), Inches(3.15), Inches(11.5), Inches(0.8))
    set_run(tt.text_frame.paragraphs[0], title, size=36, bold=True, color=WHITE)
    st = s.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.5))
    set_run(st.text_frame.paragraphs[0], subtitle, size=16, color=RGBColor(220, 240, 245))
    footer(s, n, total, "ExportIntel AI")
    return s


def main():
    ensure_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    TOTAL = 18

    # ---------- 1 Title & Team ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs, WHITE)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.35))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()
    rail = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), prs.slide_height)
    rail.fill.solid()
    rail.fill.fore_color.rgb = TEAL
    rail.line.fill.background()

    t = s.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(12), Inches(0.7))
    set_run(t.text_frame.paragraphs[0], "ExportIntel AI", size=40, bold=True, color=WHITE)
    st = s.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(0.6))
    set_run(
        st.text_frame.paragraphs[0],
        "Predictive Export Intelligence & Logistics Optimization Platform",
        size=18,
        color=RGBColor(200, 220, 235),
    )

    card(s, 0.7, 2.8, 12.0, 3.7)
    info = s.shapes.add_textbox(Inches(1.0), Inches(3.05), Inches(11.4), Inches(3.2))
    tf = info.text_frame
    set_run(tf.paragraphs[0], "Title & Team Overview", size=20, bold=True, color=NAVY)
    lines = [
        "Team Name: ____________________________    Team Number: ______",
        "Project Leader: ________________________",
        "",
        "Team Members:",
        "1) ____________________    2) ____________________",
        "3) ____________________    4) ____________________",
        "",
        "Pitch focus: Demand → Price → Logistics → Containers → Actionable Export Decision",
    ]
    for line in lines:
        p = tf.add_paragraph()
        set_run(p, line, size=16, color=INK)
        p.space_after = Pt(4)
    footer(s, 1, TOTAL)

    # ---------- 2 Problem ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Problem Statement", "Why exporters lose speed, clarity, and margin")
    bullets(
        s,
        [
            "Indian exporters must decide what to export, where to sell, and which route is profitable — under daily market shifts.",
            "Demand signals, mandi prices, freight costs, and container limits sit in separate tools and spreadsheets.",
            "Manual planning is slow, inconsistent, and often ignores live news that changes opportunity overnight.",
            "Without a unified decision system, teams risk wrong country–commodity choices and loss-making lanes.",
        ],
        0.7,
        1.7,
        11.9,
        4.7,
        size=18,
    )
    footer(s, 2, TOTAL)

    # ---------- 3 Proposed Solution ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Proposed Solution", "One multi-agent platform for export decisions")
    bullets(
        s,
        [
            "ExportIntel AI is a React + Flask product that runs a full decision pipeline in one click.",
            "News Agent pulls live demand & India price headlines (GDELT + Google News RSS).",
            "ML agents predict next-month demand opportunities and commodity prices (INR/quintal).",
            "Logistics + container agents optimize India→destination corridors and allocate scarce capacity.",
            "Groq LLM + RAG explain decisions and answer trade questions from trusted project data.",
        ],
        0.7,
        1.7,
        11.9,
        4.7,
        size=17,
    )
    footer(s, 3, TOTAL)

    # ---------- 4 Target Users ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Target Users & Use Cases")
    # three cards
    cols = [
        ("Export Analysts", ["Compare demand scores", "Check next-month buy price", "Pick profitable lanes"]),
        ("Operations Teams", ["Allocate limited containers", "Prioritize export-first lanes", "Track transit & cost/ton"]),
        ("Decision Makers", ["Review supervisor summary", "Validate profit in INR", "Ask Trade Assistant (RAG)"]),
    ]
    for i, (h, items) in enumerate(cols):
        x = 0.55 + i * 4.15
        card(s, x, 1.65, 3.95, 4.8)
        ht = s.shapes.add_textbox(Inches(x + 0.2), Inches(1.85), Inches(3.5), Inches(0.45))
        set_run(ht.text_frame.paragraphs[0], h, size=18, bold=True, color=TEAL)
        bx = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.45), Inches(3.5), Inches(3.7))
        tf = bx.text_frame
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            set_run(p, f"• {it}", size=15, color=INK)
            p.space_after = Pt(10)
    footer(s, 4, TOTAL)

    # ---------- 5 Tech Stack ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Tech Stack", "Production-style local stack with cloud scale path")
    left = [
        "Frontend: React, Vite",
        "Backend: Flask REST APIs",
        "ML: XGBoost / Joblib, Pandas",
        "LLM: Groq (Llama 3.3 70B)",
        "RAG: Chroma Vector DB (+ sklearn fallback)",
        "News: GDELT, Google News RSS, optional NewsAPI",
        "Store: SQLite (pipeline + logistics runs)",
    ]
    right = [
        "Optional scale path:",
        "Kafka (streaming news topics)",
        "Databricks + PySpark notebooks",
        "Delta bronze → silver → gold",
        "Docker for local services",
        "FX display rate: ₹96.3 / USD",
        "Auth-protected pipeline APIs",
    ]
    bullets(s, left, 0.65, 1.7, 5.8, 4.7, size=16)
    bullets(s, right, 6.85, 1.7, 5.8, 4.7, size=16)
    footer(s, 5, TOTAL)

    # ---------- 6 Data Source ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Data Source Breakdown", "Original project data + live sources + controlled synthetic support")
    bullets(
        s,
        [
            "Original / trusted datasets: Demand_prediction CSVs, commodities monthly price & training data, Logistics + Logistics_Costs CSVs (ports, freight, charges, trade prices).",
            "Live external data: GDELT + Google News RSS for demand-country and India mandi headlines (filtered to project commodities/countries).",
            "Model artifacts: demand_model_bundle.joblib and price model tools (trained on project datasets).",
            "Synthetic / generated support: synthetic news/feature rows used for training coverage and stress cases — not for final demo truth.",
            "Runtime persistence: SQLite stores pipeline JSON + logistics cost/profit rows after each analysis.",
        ],
        0.7,
        1.7,
        11.9,
        4.7,
        size=16,
    )
    footer(s, 6, TOTAL)

    # ---------- 7 Workflow ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Technical Workflow", "Multi-agent pipeline from news to export action")
    bullets(
        s,
        [
            "1) News Agent → live demand news + India price news",
            "2) Demand Agent → top country–commodity demand scores",
            "3) Price Agent → next-month INR/quintal forecasts (news-adjusted)",
            "4) Logistics Agent → India port → destination route + net profit",
            "5) Container Agent → allocate limited 20FT/40FT containers by priority",
            "6) Explain Agent (Groq) + RAG Assistant → reasoning & Q&A",
            "Output: supervisor recommendation + dashboard cards + SQLite history",
        ],
        0.7,
        1.65,
        11.9,
        4.8,
        size=16,
    )
    footer(s, 7, TOTAL)

    # ---------- 8 Architecture image ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "System Architecture", "Agents, APIs, models, and knowledge layer")
    if IMG_ARCH.exists():
        s.shapes.add_picture(str(IMG_ARCH), Inches(0.55), Inches(1.45), Inches(12.2), Inches(5.3))
    footer(s, 8, TOTAL)

    # ---------- 9 Existing vs Ours ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Existing System vs ExportIntel AI", "Impact analysis & competitive edge")
    card(s, 0.55, 1.6, 5.9, 5.0)
    card(s, 6.85, 1.6, 5.9, 5.0)
    lh = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(0.4))
    set_run(lh.text_frame.paragraphs[0], "Typical / Existing Approach", size=17, bold=True, color=ACCENT)
    rh = s.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.4), Inches(0.4))
    set_run(rh.text_frame.paragraphs[0], "Our System", size=17, bold=True, color=TEAL)
    left_items = [
        "Spreadsheet + separate tools",
        "Static news copy-paste",
        "Price/demand viewed in isolation",
        "Route cost checked later",
        "Container plan is manual",
        "Hard to explain decisions",
    ]
    right_items = [
        "One click end-to-end pipeline",
        "Auto live news (GDELT/RSS)",
        "Linked demand → price → profit",
        "Port corridors ranked by INR profit",
        "Priority-based container allocation",
        "LLM explanations + RAG Q&A",
    ]
    lb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5.4), Inches(3.9))
    rb = s.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.4), Inches(3.9))
    for box, items in ((lb, left_items), (rb, right_items)):
        tf = box.text_frame
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            set_run(p, f"• {it}", size=15, color=INK)
            p.space_after = Pt(9)
    footer(s, 9, TOTAL)

    # ---------- 10 Demand & Price deep dive ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Core Intelligence: Demand & Price", "How scores and forecasts are produced")
    bullets(
        s,
        [
            "Demand score (0–1): news features (shortage, sentiment, opportunity) + trained demand model → ranked opportunities.",
            "Price forecast: dataset current mandi baseline + XGBoost next-month prediction + explicit news sentiment adjustment.",
            "Same news → stable prediction; clearly bullish/bearish news → price moves (verified in product behavior).",
            "All money values for logistics/profit shown in INR (FX ₹96.3/USD); prices stay INR/quintal.",
        ],
        0.7,
        1.7,
        11.9,
        4.7,
        size=16,
    )
    footer(s, 10, TOTAL)

    # ---------- 11 Logistics profit ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Logistics & Profit Engine", "Buying cost + quantity-aware allocation")
    bullets(
        s,
        [
            "Net profit / ton = Sell (trade price) − Buy (predicted India price) − Logistics cost / ton.",
            "Buying amount is always considered via predicted INR/quintal converted to per-ton cost.",
            "Quantity enters at container & allocation level: payload tons × containers allocated.",
            "Logistics page compares routes on per-ton profit; Container page shows combined expected profit.",
        ],
        0.7,
        1.7,
        11.9,
        4.7,
        size=16,
    )
    footer(s, 11, TOTAL)

    # ---------- 12 Performance & Demo ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "System Performance & Live Demo", "Local deployment ready for review")
    bullets(
        s,
        [
            "Frontend: http://localhost:5173  |  Backend Flask API: http://127.0.0.1:5001",
            "Live demo path: Login → Fetch Live News / Run Analysis → Demand → Price → Logistics → Containers → Agents/RAG.",
            "SQLite retains pipeline runs and logistics cost rows for auditability.",
            "RAG answers use retrieved chunks from demand/price/logistics datasets (Chroma).",
            "Pitch demo highlight: one analysis produces ranked demand, INR profit routes, and container plan.",
        ],
        0.7,
        1.65,
        11.9,
        4.8,
        size=16,
    )
    footer(s, 12, TOTAL)

    # ---------- 13 Demo screenshots ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Live Product Screens", "Dashboard · Demand · Logistics · Containers")
    imgs = [IMG3, IMG2, IMG1, IMG4]
    labels = ["Dashboard", "Demand", "Logistics", "Containers"]
    coords = [(0.5, 1.55), (6.7, 1.55), (0.5, 4.15), (6.7, 4.15)]
    for img, (x, y), lab in zip(imgs, coords, labels):
        cap = s.shapes.add_textbox(Inches(x), Inches(y - 0.22), Inches(2.5), Inches(0.22))
        set_run(cap.text_frame.paragraphs[0], lab, size=11, bold=True, color=MUTED)
        if img.exists():
            s.shapes.add_picture(str(img), Inches(x), Inches(y), Inches(5.95), Inches(2.25))
    footer(s, 13, TOTAL)

    # ---------- 14 Alternative Approach ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Alternative Approach Considered", "Why we selected the primary design")
    bullets(
        s,
        [
            "Alternative A: Pure LLM decisioning without ML models — rejected (less stable, harder to audit).",
            "Alternative B: Manual pasted news only — rejected (slow, not product-like for daily decisions).",
            "Alternative C: Databricks + Kafka-first for every demo — deferred (great for scale; overkill for local review).",
            "Selected approach: Hybrid ML + live news APIs + Flask/React product UX, with Kafka/Databricks as scale roadmap.",
            "Result: reliable demos, explainable scores, and a clear path to streaming production.",
        ],
        0.7,
        1.65,
        11.9,
        4.8,
        size=16,
    )
    footer(s, 14, TOTAL)

    # ---------- 15 Future Scope ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Future Scope & Enhancements", "Roadmap beyond the review demo")
    bullets(
        s,
        [
            "Always-on Kafka news producers → Databricks bronze/silver/gold tables.",
            "Dynamic FX feeds and confidence intervals on price/demand predictions.",
            "What-if simulator for container count, port constraints, and tariff shocks.",
            "Role-based multi-user workspaces and exportable board packs (PDF/PPT).",
            "Stronger evaluation harness: backtests on historical news weeks + route outcomes.",
        ],
        0.7,
        1.65,
        11.9,
        4.8,
        size=16,
    )
    footer(s, 15, TOTAL)

    # ---------- 16 Impact snapshot ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Product Impact Snapshot", "What the pitch delivers in 14 minutes")
    bullets(
        s,
        [
            "Faster export decisions: one pipeline replaces fragmented analysis.",
            "Profit-first logistics: buy price, freight, and route compared in INR.",
            "Capacity realism: limited containers allocated by opportunity priority.",
            "Trust layer: LLM explanations + RAG grounded in project datasets.",
            "Demo-ready: live news → ranked opportunities → actionable recommendation.",
        ],
        0.7,
        1.7,
        11.9,
        4.7,
        size=17,
    )
    footer(s, 16, TOTAL)

    # ---------- 17 Conclusion ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs, NAVY)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), prs.slide_width, Inches(3.1))
    strip.fill.solid()
    strip.fill.fore_color.rgb = TEAL
    strip.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(11.7), Inches(0.8))
    set_run(t.text_frame.paragraphs[0], "Conclusion", size=34, bold=True, color=WHITE)
    st = s.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(11.7), Inches(1.4))
    set_run(
        st.text_frame.paragraphs[0],
        "ExportIntel AI turns live market news into ranked demand, price forecasts,\nprofitable routes, and container plans — with explainable AI support.",
        size=18,
        color=WHITE,
    )
    q = s.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5))
    set_run(q.text_frame.paragraphs[0], "Q&A Ready  ·  Thank you", size=20, bold=True, color=RGBColor(220, 240, 245))
    footer(s, 17, TOTAL, "ExportIntel AI")

    # ---------- 18 Backup / speaking split ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, prs)
    title_block(s, "Suggested Speaking Split (~14 min)", "Equal participation plan for team members")
    bullets(
        s,
        [
            "Member A (3 min): Title, problem, proposed solution, users",
            "Member B (3 min): Tech stack, data sources, architecture/workflow",
            "Member C (4 min): Demand/price logic, logistics profit, existing vs ours",
            "Member D (3 min): Live demo screens, alternatives, future scope, close/Q&A",
            "Tip: Keep demo to 1 Run Analysis path — Demand card → Route → Containers.",
        ],
        0.7,
        1.7,
        11.9,
        4.7,
        size=16,
    )
    footer(s, 18, TOTAL)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
