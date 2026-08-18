"""
ExportIntel AI — Clean project presentation
Title uses Platform cover image; four phases each get a detailed slide.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "ExportIntel_AI_Project.pptx"
OUT_DL = Path(r"c:\Users\Lenovo\Downloads\ExportIntel_AI_Project.pptx")

ASSETS = ROOT / "assets" / "ppt"
TITLE_IMG = ASSETS / "title_platform.png"
TITLE_FALLBACK = ASSETS / "platform_slides" / "slide_01.png"
ARCH = ASSETS / "architecture_chroma.png"
ARCH_FALLBACK = ROOT / "ExportIntel_AI_Architecture_Diagram.png"
DEMO_D = ASSETS / "demo_dashboard.png"
DEMO_DEM = ASSETS / "demo_demand.png"
DEMO_L = ASSETS / "demo_logistics.png"
DEMO_C = ASSETS / "demo_containers.png"
SHIP = ASSETS / "ship_hero.png"
PORT = ASSETS / "port_terminal.png"
CONT = ASSETS / "containers_yard.png"
OCEAN = ASSETS / "ocean_trade.png"

NAVY = RGBColor(8, 18, 40)
NAVY2 = RGBColor(12, 28, 58)
PANEL = RGBColor(16, 36, 68)
INK = RGBColor(230, 236, 245)
MUTED = RGBColor(160, 175, 195)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(30, 190, 200)
GOLD = RGBColor(230, 170, 55)
CORAL = RGBColor(230, 90, 70)
SOFT = RGBColor(50, 80, 120)

SLIDE_W = 13.333
TOTAL = 17


def run(p, text, *, size=16, bold=False, color=INK, italic=False, align=None):
    p.clear()
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if align is not None:
        p.alignment = align


def dark_bg(slide, prs, photo=None, dim=0.78):
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    base.fill.solid()
    base.fill.fore_color.rgb = NAVY
    base.line.fill.background()
    if photo and Path(photo).exists():
        slide.shapes.add_picture(str(photo), 0, 0, prs.slide_width, prs.slide_height)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = NAVY
        overlay.line.fill.background()
        try:
            from lxml import etree

            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            solidFill = overlay._element.spPr.find(f"{ns}solidFill")
            srgb = solidFill.find(f"{ns}srgbClr") if solidFill is not None else None
            if srgb is not None:
                alpha = etree.SubElement(srgb, f"{ns}alpha")
                alpha.set("val", str(int(dim * 100000)))
        except Exception:
            pass
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL
    top.line.fill.background()


def glass(slide, x, y, w, h, *, fill=PANEL, line=SOFT):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1.25)
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh


def accent_bar(slide, x, y, h, color=TEAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def footer(slide, n):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.28))
    run(
        box.text_frame.paragraphs[0],
        f"ExportIntel AI  ·  Predictive Export Intelligence                              {n} / {TOTAL}",
        size=11,
        color=MUTED,
    )


def heading(slide, title, subtitle=""):
    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.55))
    run(t.text_frame.paragraphs[0], title, size=26, bold=True, color=WHITE)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(0.85), Inches(12.2), Inches(0.35))
        run(s.text_frame.paragraphs[0], subtitle, size=14, color=TEAL, italic=True)


def bullets(slide, items, x, y, w, h, *, size=15, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run(p, f"•  {item}", size=size, color=color)
        p.space_after = Pt(8)


def content(prs, title, subtitle, items, n, *, photo=None, size=16, accent=TEAL):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, photo, dim=0.72)
    heading(s, title, subtitle)
    glass(s, 0.5, 1.35, 12.3, 5.4)
    accent_bar(s, 0.5, 1.35, 5.4, accent)
    bullets(s, items, 0.85, 1.6, 11.7, 4.9, size=size)
    footer(s, n)


def paragraph_slide(prs, title, subtitle, paragraphs, n, *, photo=None, size=17, accent=CORAL):
    """Problem / narrative slide — flowing paragraphs, not bullets."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, photo, dim=0.72)
    heading(s, title, subtitle)
    glass(s, 0.5, 1.35, 12.3, 5.4)
    accent_bar(s, 0.5, 1.35, 5.4, accent)
    box = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.7))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for para in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run(p, para, size=size, color=INK)
        p.space_after = Pt(16)
        p.line_spacing = 1.25
    footer(s, n)


def team_slide(prs, n):
    members = [
        "Mohasin Nadaf",
        "Soham Kulkarni",
        "Nandini Khandade",
        "Abhishek Ingalkar",
        "Shubham Kamble",
    ]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Solid dark slide — no photo background
    dark_bg(s, prs, photo=None)

    heading(s, "Our Team", "ExportIntel AI — Predictive Export Intelligence")

    # Team name banner
    glass(s, 0.55, 1.4, 12.2, 1.15, line=TEAL)
    accent_bar(s, 0.55, 1.4, 1.15, TEAL)
    tn = s.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(3.0), Inches(0.35))
    run(tn.text_frame.paragraphs[0], "TEAM NAME", size=12, bold=True, color=MUTED)
    tv = s.shapes.add_textbox(Inches(0.9), Inches(1.95), Inches(11.4), Inches(0.45))
    run(tv.text_frame.paragraphs[0], "Solo Leveling", size=28, bold=True, color=WHITE)

    # Member cards — clean row, no photo clutter
    gap, margin = 0.28, 0.55
    usable = SLIDE_W - 2 * margin - gap * 4
    w = usable / 5
    for i, name in enumerate(members):
        x = margin + i * (w + gap)
        col = GOLD if i % 2 == 0 else TEAL
        glass(s, x, 2.95, w, 3.7, line=col)
        accent_bar(s, x, 2.95, 3.7, col)

        # Circle-like number badge
        badge = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + w / 2 - 0.35), Inches(3.35), Inches(0.7), Inches(0.7)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = col
        badge.line.fill.background()
        num = s.shapes.add_textbox(Inches(x + w / 2 - 0.35), Inches(3.48), Inches(0.7), Inches(0.45))
        run(num.text_frame.paragraphs[0], str(i + 1), size=18, bold=True, color=NAVY)
        num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        label = s.shapes.add_textbox(Inches(x + 0.15), Inches(4.3), Inches(w - 0.3), Inches(0.35))
        run(label.text_frame.paragraphs[0], "Team member", size=11, color=MUTED)
        label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        nm = s.shapes.add_textbox(Inches(x + 0.12), Inches(4.75), Inches(w - 0.24), Inches(1.4))
        ntf = nm.text_frame
        ntf.word_wrap = True
        run(ntf.paragraphs[0], name, size=15, bold=True, color=WHITE)
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer(s, n)


def two_col(prs, title, subtitle, lt, li, rt, ri, n, *, photo=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, photo, dim=0.72)
    heading(s, title, subtitle)
    glass(s, 0.5, 1.35, 6.0, 5.4)
    accent_bar(s, 0.5, 1.35, 5.4, TEAL)
    h1 = s.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(5.4), Inches(0.4))
    run(h1.text_frame.paragraphs[0], lt, size=16, bold=True, color=TEAL)
    bullets(s, li, 0.8, 2.1, 5.4, 4.3, size=14)

    glass(s, 6.8, 1.35, 6.0, 5.4)
    accent_bar(s, 6.8, 1.35, 5.4, GOLD)
    h2 = s.shapes.add_textbox(Inches(7.1), Inches(1.55), Inches(5.4), Inches(0.4))
    run(h2.text_frame.paragraphs[0], rt, size=16, bold=True, color=GOLD)
    bullets(s, ri, 7.1, 2.1, 5.4, 4.3, size=14)
    footer(s, n)


def title_slide(prs):
    """Full-bleed cover — user's Platform title image (no generated overlay)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    img = TITLE_IMG if TITLE_IMG.exists() else TITLE_FALLBACK
    dark_bg(s, prs)
    if img.exists():
        s.shapes.add_picture(str(img), 0, 0, prs.slide_width, prs.slide_height)


def phase_slide(prs, n, *, phase, title, subtitle, accent, how_it_works, inputs, outputs, photo=None, demo=None):
    """One dedicated slide per product phase with extra project detail."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, photo, dim=0.74)

    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.28), Inches(1.6), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    try:
        badge.adjustments[0] = 0.2
    except Exception:
        pass
    bt = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(1.6), Inches(0.35))
    run(bt.text_frame.paragraphs[0], f"PHASE {phase}", size=12, bold=True, color=NAVY)
    bt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    ht = s.shapes.add_textbox(Inches(2.35), Inches(0.22), Inches(10.4), Inches(0.5))
    run(ht.text_frame.paragraphs[0], title, size=26, bold=True, color=WHITE)
    st = s.shapes.add_textbox(Inches(2.35), Inches(0.75), Inches(10.4), Inches(0.35))
    run(st.text_frame.paragraphs[0], subtitle, size=13, color=TEAL, italic=True)

    # How it works — full width
    glass(s, 0.5, 1.25, 8.3, 3.35, line=accent)
    accent_bar(s, 0.5, 1.25, 3.35, accent)
    h1 = s.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(7.8), Inches(0.35))
    run(h1.text_frame.paragraphs[0], "How it works in ExportIntel AI", size=15, bold=True, color=accent)
    bullets(s, how_it_works, 0.75, 1.9, 7.8, 2.5, size=14)

    # Demo / visual card on right
    glass(s, 9.0, 1.25, 3.8, 3.35, line=accent)
    if demo and Path(demo).exists():
        lab = s.shapes.add_textbox(Inches(9.15), Inches(1.35), Inches(3.5), Inches(0.3))
        run(lab.text_frame.paragraphs[0], "Module preview", size=11, bold=True, color=MUTED)
        s.shapes.add_picture(str(demo), Inches(9.2), Inches(1.75), Inches(3.4), Inches(2.6))
    else:
        lab = s.shapes.add_textbox(Inches(9.2), Inches(2.4), Inches(3.4), Inches(1.0))
        run(lab.text_frame.paragraphs[0], f"Dashboard\n{title}", size=16, bold=True, color=WHITE)
        lab.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Inputs
    glass(s, 0.5, 4.8, 6.0, 2.0)
    accent_bar(s, 0.5, 4.8, 2.0, GOLD)
    ih = s.shapes.add_textbox(Inches(0.75), Inches(4.95), Inches(5.5), Inches(0.3))
    run(ih.text_frame.paragraphs[0], "Inputs", size=14, bold=True, color=GOLD)
    bullets(s, inputs, 0.75, 5.35, 5.5, 1.3, size=13)

    # Outputs
    glass(s, 6.8, 4.8, 6.0, 2.0)
    accent_bar(s, 6.8, 4.8, 2.0, TEAL)
    oh = s.shapes.add_textbox(Inches(7.05), Inches(4.95), Inches(5.5), Inches(0.3))
    run(oh.text_frame.paragraphs[0], "Outputs", size=14, bold=True, color=TEAL)
    bullets(s, outputs, 7.05, 5.35, 5.5, 1.3, size=13)

    footer(s, n)


def pipeline_slide(prs, n):
    steps = [
        ("News", "GDELT + Google RSS", MUTED),
        ("Demand", "Opportunity scores", TEAL),
        ("Price", "INR/quintal forecast", GOLD),
        ("Logistics", "Net profit routes", CORAL),
        ("Containers", "Priority allocation", TEAL),
        ("Explain", "Groq + RAG", MUTED),
    ]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, CONT if CONT.exists() else None, dim=0.78)
    heading(s, "Run Analysis Pipeline", "Real backend order — one click, full decision chain")

    gap, margin = 0.2, 0.45
    usable = SLIDE_W - 2 * margin - gap * 5
    w = usable / 6
    for i, (title, blurb, col) in enumerate(steps):
        x = margin + i * (w + gap)
        glass(s, x, 1.6, w, 3.8, line=col)
        accent_bar(s, x, 1.6, 3.8, col)
        num = s.shapes.add_textbox(Inches(x + 0.15), Inches(1.85), Inches(w - 0.25), Inches(0.35))
        run(num.text_frame.paragraphs[0], f"{i + 1}", size=18, bold=True, color=col)
        ht = s.shapes.add_textbox(Inches(x + 0.15), Inches(2.45), Inches(w - 0.25), Inches(0.9))
        run(ht.text_frame.paragraphs[0], title, size=15, bold=True, color=WHITE)
        bt = s.shapes.add_textbox(Inches(x + 0.15), Inches(3.5), Inches(w - 0.25), Inches(1.4))
        btf = bt.text_frame
        btf.word_wrap = True
        run(btf.paragraphs[0], blurb, size=12, color=MUTED)
        if i < 5:
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(x + w + 0.01), Inches(3.35), Inches(0.18), Inches(0.2)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = SOFT
            arr.line.fill.background()

    note = s.shapes.add_textbox(Inches(0.55), Inches(5.7), Inches(12.2), Inches(1.0))
    ntf = note.text_frame
    ntf.word_wrap = True
    run(
        ntf.paragraphs[0],
        "UI modules: Dashboard · Price Prediction · Demand Prediction · Logistics Optimisation · "
        "Container Priority · AI Assistant · Agent Reasoning · Analytics",
        size=13,
        color=INK,
    )
    footer(s, n)


def arch_slide(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs)
    heading(s, "System Architecture", "Frontend · Flask APIs · ML models · RAG · SQLite")
    img = ARCH if ARCH.exists() else ARCH_FALLBACK
    if img.exists():
        s.shapes.add_picture(str(img), Inches(0.55), Inches(1.3), Inches(12.2), Inches(5.5))
    footer(s, n)


def demo_grid(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs)
    heading(s, "Live Product Screens", "Real ExportIntel AI dashboard modules")
    tiles = [
        (DEMO_D, "Dashboard"),
        (DEMO_DEM, "Demand Prediction"),
        (DEMO_L, "Logistics Optimisation"),
        (DEMO_C, "Container Priority"),
    ]
    coords = [(0.45, 1.4), (6.85, 1.4), (0.45, 4.2), (6.85, 4.2)]
    for (img, label), (x, y) in zip(tiles, coords):
        glass(s, x, y, 6.0, 2.55)
        lab = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(5.5), Inches(0.3))
        run(lab.text_frame.paragraphs[0], label, size=12, bold=True, color=TEAL)
        if Path(img).exists():
            s.shapes.add_picture(str(img), Inches(x + 0.25), Inches(y + 0.45), Inches(5.5), Inches(1.9))
    footer(s, n)


def thank_you(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, OCEAN if OCEAN.exists() else None, dim=0.7)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(2.0), Inches(9.3), Inches(3.4))
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY2
    panel.line.color.rgb = TEAL
    try:
        panel.adjustments[0] = 0.05
    except Exception:
        pass

    t = s.shapes.add_textbox(Inches(2.3), Inches(2.35), Inches(8.7), Inches(0.7))
    run(t.text_frame.paragraphs[0], "Thank You", size=40, bold=True, color=WHITE)
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    st = s.shapes.add_textbox(Inches(2.3), Inches(3.15), Inches(8.7), Inches(0.45))
    run(st.text_frame.paragraphs[0], "Predictive Export Intelligence", size=20, color=TEAL)
    st.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    sub = s.shapes.add_textbox(Inches(2.3), Inches(3.8), Inches(8.7), Inches(1.0))
    run(
        sub.text_frame.paragraphs[0],
        "Price  →  Demand  →  Logistics  →  Containers\nQuestions & discussion",
        size=15,
        color=MUTED,
    )
    sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer(s, n)


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(7.5)

    # 1 — Title cover
    title_slide(prs)

    # 2 — Team
    team_slide(prs, 2)

    # 3 Agenda
    content(
        prs,
        "Agenda",
        "Project walkthrough",
        [
            "1. Team introduction and problem statement",
            "2. Project objectives, architecture, and Run Analysis pipeline",
            "3. Four modules — Price, Demand, Logistics, Containers",
            "4. Tech stack, live screens, explainability, and results",
            "5. Future scope and close",
        ],
        3,
        photo=PORT if PORT.exists() else None,
        size=18,
    )

    # 4 Problem — paragraph form
    paragraph_slide(
        prs,
        "Problem Statement",
        "Why Indian commodity exporters need a unified system",
        [
            "Indian commodity exporters operate in fast-shifting global markets, yet the data they need for daily decisions is trapped across scattered tools. Demand signals, live market news, mandi prices, freight costs, and container limits rarely sit in one place, so teams spend hours stitching spreadsheets instead of acting on opportunity.",
            "Manual planning is slow and often misses overnight news shifts that change country–commodity attractiveness or India buy prices. Route selection and container allocation — the choices that decide real profit — are frequently made late, ad-hoc, or without a clear net-profit view.",
            "Without a unified predictive pipeline, exporters risk wrong country–commodity choices and loss-making logistics lanes. ExportIntel AI addresses this by linking price prediction, demand scoring, logistics optimization, and container priority into one explainable decision platform.",
        ],
        4,
        photo=OCEAN if OCEAN.exists() else None,
        size=16,
        accent=CORAL,
    )

    # 5 Objectives
    content(
        prs,
        "Project Objectives",
        "ExportIntel AI — Predictive Export Intelligence",
        [
            "1. Commodity price prediction — next-month India buy prices (INR/quintal)",
            "2. Demand prediction — top country–commodity opportunities from live news + ML",
            "3. Logistic optimization — India port → destination corridors by INR net profit",
            "4. Container priority — allocate limited 20FT/40FT capacity by opportunity score",
            "Explainable outputs via Groq LLM + RAG Trade Assistant",
        ],
        5,
        photo=CONT if CONT.exists() else None,
        size=16,
        accent=GOLD,
    )

    # 6 Architecture (before price / modules)
    arch_slide(prs, 6)

    # 7 Pipeline (right after architecture)
    pipeline_slide(prs, 7)

    # 8 — Phase 1 Price
    phase_slide(
        prs,
        8,
        phase="01",
        title="Commodity Price Prediction",
        subtitle="Next-month India buy price in INR/quintal — Price Prediction module",
        accent=GOLD,
        how_it_works=[
            "Reads latest mandi baseline from commodities monthly_price.csv (e.g. 2026-06)",
            "XGBoost model (commodities/price_agent_tools) forecasts next-month price",
            "India market news sentiment applies a controlled +/- adjustment",
            "Supported commodities: Coffee, Cotton, Maize, Onion, Soybean, Sugar, Turmeric, Wheat",
            "Groq Explain Agent narrates why the price moved for Agent Reasoning",
        ],
        inputs=[
            "Mandi / monthly price history (INR/quintal)",
            "Live India commodity news (GDELT + Google RSS)",
            "MA7 / MA30 and recent price-change features",
        ],
        outputs=[
            "Predicted next-month buy price (INR/quintal)",
            "Sentiment-adjusted forecast vs baseline",
            "Feeds logistics buy-cost and profit math",
        ],
        photo=PORT if PORT.exists() else None,
        demo=DEMO_D if DEMO_D.exists() else None,
    )

    # 9 — Phase 2 Demand
    phase_slide(
        prs,
        9,
        phase="02",
        title="Demand Prediction",
        subtitle="Country–commodity opportunity scoring — Demand Prediction module",
        accent=TEAL,
        how_it_works=[
            "News Agent pulls demand-country headlines filtered to project commodities",
            "Headlines parsed for country, commodity, shortage, and sentiment signals",
            "Trained demand_model_bundle.joblib scores each pair from 0 to 1",
            "Shortage / production drop / export opportunity raise the demand score",
            "Top ranked opportunities become inputs to logistics and container stages",
        ],
        inputs=[
            "Live demand news (target countries + commodities)",
            "Demand prediction datasets + trained Joblib model",
            "Extracted news features (sentiment, shortage flags)",
        ],
        outputs=[
            "Ranked country–commodity demand scores (0–1)",
            "Top export opportunity list on dashboard",
            "Demand score used in lane profit & priority",
        ],
        photo=OCEAN if OCEAN.exists() else None,
        demo=DEMO_DEM if DEMO_DEM.exists() else None,
    )

    # 10 — Phase 3 Logistics
    phase_slide(
        prs,
        10,
        phase="03",
        title="Logistic Optimization",
        subtitle="India port → destination corridors ranked by INR net profit",
        accent=CORAL,
        how_it_works=[
            "Builds lanes from India origin ports to destination countries for each opportunity",
            "Net profit/ton = sell (trade price) − buy (predicted India price) − logistics cost/ton",
            "Logistics_Costs CSVs supply freight, port charges, and transit data",
            "All money shown in INR at fixed FX: 1 USD = ₹96.3",
            "Pipeline keeps only profitable lanes before container allocation",
        ],
        inputs=[
            "Predicted India buy price (INR/quintal)",
            "Demand score + trade sell prices",
            "Port / freight / charges CSVs (Logistics_Costs)",
        ],
        outputs=[
            "Corridors ranked by INR net profit/ton",
            "Cost breakdown: buy, sell, logistics",
            "Profitable lane set for container stage",
        ],
        photo=CONT if CONT.exists() else None,
        demo=DEMO_L if DEMO_L.exists() else None,
    )

    # 11 — Phase 4 Containers
    phase_slide(
        prs,
        11,
        phase="04",
        title="Container Priority",
        subtitle="Scarce capacity allocated across competing export lanes",
        accent=TEAL,
        how_it_works=[
            "Default planning example: limited fleet (e.g. 6 × 20FT) vs many opportunities",
            "Container_prioritization ranks lanes by demand + net profit + cost + transit",
            "Higher-priority lanes get more containers; lower ones get fewer or zero",
            "Quantity applied here: profit/container = profit/ton × payload tons",
            "Dashboard shows export-first allocation plan for operations teams",
        ],
        inputs=[
            "Profitable logistics lanes + demand scores",
            "Available container count and type (20FT/40FT)",
            "Transit time and logistics cost per lane",
        ],
        outputs=[
            "Priority-ranked container allocation plan",
            "Containers per lane + projected profit",
            "Actionable export recommendation on UI",
        ],
        photo=SHIP if SHIP.exists() else None,
        demo=DEMO_C if DEMO_C.exists() else None,
    )

    # 12 Tech stack
    two_col(
        prs,
        "Technology Stack",
        "What the project is built with",
        "Application",
        [
            "Frontend: React 18 + Vite",
            "Backend: Flask REST APIs (:5001)",
            "Data: Pandas / NumPy",
            "DB: SQLite (pipeline + logistics history)",
            "Auth: demo login for protected APIs",
        ],
        "AI & Data",
        [
            "ML: XGBoost + Joblib (price & demand)",
            "LLM: Groq Llama 3.3 70B",
            "RAG: Chroma Vector DB",
            "News: GDELT + Google News RSS",
            "Explain Agent for reasoning text",
        ],
        12,
        photo=PORT if PORT.exists() else None,
    )

    # 13 Demo
    demo_grid(prs, 13)

    # 14 Value
    two_col(
        prs,
        "Explainability & Business Value",
        "Grounded answers + faster export decisions",
        "RAG + Agent Reasoning",
        [
            "Trusted CSVs indexed in Chroma Vector DB",
            "AI Assistant answers from retrieved context",
            "Groq Llama 3.3 70B grounds answers in chunks",
            "Agent Reasoning explains demand & price calls",
        ],
        "Business Value",
        [
            "One-click Run Analysis replaces fragmented tools",
            "Profit visibility before shipment commitment",
            "Scarce containers allocated by real priority",
            "Clear, explainable output for leadership",
        ],
        14,
        photo=PORT if PORT.exists() else None,
    )

    # 15 Results
    content(
        prs,
        "Results & Demo Outcomes",
        "What the working product delivers",
        [
            "Linked decisions across all four modules on one dashboard",
            "Live news-aware price and demand predictions",
            "Corridors ranked by INR net profit before commit",
            "Practical limited-capacity container plans",
            "SQLite stores pipeline runs for audit / replay",
        ],
        15,
        photo=SHIP if SHIP.exists() else None,
        size=17,
    )

    # 16 Future scope (before thank you)
    content(
        prs,
        "Future Scope",
        "How ExportIntel AI can grow beyond the current demo",
        [
            "Kafka-based streaming so market news and prices update continuously",
            "Databricks-scale batch/stream pipelines for larger commodity and route datasets",
            "Dynamic FX feeds instead of fixed ₹96.3/USD conversion",
            "Confidence intervals and risk bands on price and demand predictions",
            "What-if simulator for container count, port constraints, and tariff shocks",
            "Mobile / ops alerts when a high-priority lane or news event appears",
        ],
        16,
        photo=OCEAN if OCEAN.exists() else None,
        size=16,
        accent=TEAL,
    )

    # 17 Thank you
    thank_you(prs, 17)

    for path in (OUT, OUT_DL):
        try:
            prs.save(str(path))
            print(f"Saved: {path}")
        except PermissionError:
            alt = path.with_name(path.stem + "_new" + path.suffix)
            prs.save(str(alt))
            print(f"Locked — saved: {alt}")

    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
