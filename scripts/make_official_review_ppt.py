"""
ExportIntel AI — Official Review PPT
Follows event structure exactly (15–20 slides). Custom navy/teal design.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "ExportIntel_AI_Official_Review_PPT.pptx"

IMG_ARCH = ROOT / "ExportIntel_AI_Architecture_Diagram.png"
IMG1 = ROOT / "assets" / "image-1.png"
IMG2 = ROOT / "assets" / "image-2.png"
IMG3 = ROOT / "assets" / "image-3.png"
IMG4 = ROOT / "assets" / "image-4.png"

NAVY = RGBColor(12, 32, 64)
TEAL = RGBColor(8, 110, 130)
INK = RGBColor(25, 32, 42)
MUTED = RGBColor(95, 105, 118)
SOFT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(214, 223, 232)
GOLD = RGBColor(180, 110, 40)

SOURCE_ASSETS = [
    ROOT.parent / ".cursor" / "projects" / "c-Users-Lenovo-Desktop-Export-AI" / "assets" / "c__Users_Lenovo_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-326990d6-d86f-4625-aad4-ffff9f82faaa.png",
    ROOT.parent / ".cursor" / "projects" / "c-Users-Lenovo-Desktop-Export-AI" / "assets" / "c__Users_Lenovo_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-8ba07128-fbb0-4042-8280-29648355cab1.png",
    ROOT.parent / ".cursor" / "projects" / "c-Users-Lenovo-Desktop-Export-AI" / "assets" / "c__Users_Lenovo_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-1e0864a7-b10a-42aa-bd69-a41c45320d26.png",
    ROOT.parent / ".cursor" / "projects" / "c-Users-Lenovo-Desktop-Export-AI" / "assets" / "c__Users_Lenovo_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-d82a3a01-86ed-4cf0-bb52-7b8d5ae2353d.png",
]


def ensure_assets():
    (ROOT / "assets").mkdir(exist_ok=True)
    for src, dst in zip(SOURCE_ASSETS, [IMG1, IMG2, IMG3, IMG4]):
        if src.exists() and not dst.exists():
            dst.write_bytes(src.read_bytes())


def run(p, text, *, size=17, bold=False, color=INK):
    p.clear()
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def bg(slide, prs, color=SOFT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
    rail.fill.solid()
    rail.fill.fore_color.rgb = TEAL
    rail.line.fill.background()


def footer(slide, n, total):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.28))
    run(box.text_frame.paragraphs[0], f"ExportIntel AI  ·  Startup Pitch Review                              {n} / {total}", size=10, color=MUTED)


def heading(slide, title, subtitle=""):
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.65))
    run(t.text_frame.paragraphs[0], title, size=28, bold=True, color=NAVY)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(0.92), Inches(12.3), Inches(0.35))
        run(s.text_frame.paragraphs[0], subtitle, size=14, color=MUTED)


def card(slide, x, y, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = LINE
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh


def bullets(slide, items, x, y, w, h, *, size=16):
    card(slide, x - 0.05, y - 0.08, w + 0.1, h + 0.16)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run(p, f"•  {item}", size=size, color=INK)
        p.space_after = Pt(7)


def main():
    ensure_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    TOTAL = 17

    # 1. Title & Team Overview
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs, WHITE)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.2))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()
    rail = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
    rail.fill.solid()
    rail.fill.fore_color.rgb = TEAL
    rail.line.fill.background()

    t = s.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12), Inches(0.7))
    run(t.text_frame.paragraphs[0], "ExportIntel AI", size=38, bold=True, color=WHITE)
    st = s.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.55))
    run(st.text_frame.paragraphs[0], "Predictive Export Intelligence & Logistics Optimization Platform", size=17, color=RGBColor(205, 225, 235))

    card(s, 0.55, 2.6, 12.2, 4.0)
    info = s.shapes.add_textbox(Inches(0.85), Inches(2.85), Inches(11.6), Inches(3.5))
    tf = info.text_frame
    run(tf.paragraphs[0], "Title & Team Overview", size=18, bold=True, color=NAVY)
    for line in [
        "",
        "Team Number: __________     Team Name: _______________________________",
        "Project Leader: _______________________________________________________",
        "",
        "Team Members (names on this slide as required):",
        "1) ________________________     2) ________________________",
        "3) ________________________     4) ________________________",
        "",
        "Product pitch: Demand → Price → Logistics → Containers → Export Decision",
    ]:
        p = tf.add_paragraph()
        run(p, line, size=15, color=INK)
        p.space_after = Pt(3)
    footer(s, 1, TOTAL)

    # 2. Problem Statement
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Problem Statement", "Core problem we are solving")
    bullets(
        s,
        [
            "Indian commodity exporters must decide what to export, which country to target, and which route is profitable — under fast-changing markets.",
            "Demand signals, mandi prices, freight costs, and container limits are scattered across tools and spreadsheets.",
            "Manual planning is slow and often misses live news that shifts opportunity overnight.",
            "Without a unified system, teams risk wrong country–commodity choices and loss-making logistics lanes.",
        ],
        0.6,
        1.55,
        12.1,
        4.9,
        size=17,
    )
    footer(s, 2, TOTAL)

    # 3. Proposed Solution
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Proposed Solution", "High-level overview of our unique product")
    bullets(
        s,
        [
            "ExportIntel AI is a multi-agent export decision platform (React + Flask).",
            "One click runs: Live News → Demand → Price → Logistics → Container Prioritization.",
            "Live news from GDELT + Google News RSS feeds demand and India price intelligence.",
            "ML models score opportunities and forecast next-month prices in INR/quintal.",
            "Logistics engine ranks India→destination corridors by net profit (INR @ ₹96.3/USD).",
            "Container Agent allocates limited capacity by demand + profit + cost + transit priority.",
            "Groq LLM explanations + RAG Trade Assistant make decisions explainable and auditable.",
        ],
        0.6,
        1.5,
        12.1,
        5.0,
        size=16,
    )
    footer(s, 3, TOTAL)

    # 4. Target Users & Use Cases
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Target Users & Use Cases", "Who uses it and where it creates value")
    cols = [
        ("Export Analysts", ["Rank demand opportunities", "Review next-month buy prices", "Compare profitable corridors"]),
        ("Operations Teams", ["Allocate limited containers", "Prioritize export-first lanes", "Track cost, transit, profit"]),
        ("Decision Makers", ["Use supervisor recommendation", "Validate INR profit impact", "Ask Trade Assistant (RAG)"]),
    ]
    for i, (h, items) in enumerate(cols):
        x = 0.5 + i * 4.2
        card(s, x, 1.5, 4.0, 5.0)
        ht = s.shapes.add_textbox(Inches(x + 0.2), Inches(1.7), Inches(3.6), Inches(0.45))
        run(ht.text_frame.paragraphs[0], h, size=17, bold=True, color=TEAL)
        bx = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.3), Inches(3.6), Inches(3.9))
        tf = bx.text_frame
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run(p, f"• {it}", size=15, color=INK)
            p.space_after = Pt(12)
    footer(s, 4, TOTAL)

    # 5. Tech Stack
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Tech Stack", "Technical architecture & technologies used")
    bullets(
        s,
        [
            "Frontend: React + Vite dashboard (Demand, Price, Logistics, Containers, Agents, RAG)",
            "Backend: Flask REST APIs (pipeline, live news, RAG ask/rebuild, auth)",
            "ML: XGBoost/Joblib demand & price models with Pandas feature pipelines",
            "LLM: Groq Llama 3.3 70B for demand/price explanations and RAG answers",
            "Vector DB: Chroma (sklearn TF-IDF fallback) for retrieval-augmented Q&A",
            "Live News: GDELT Doc API + Google News RSS (+ optional NewsAPI)",
            "Persistence: SQLite for pipeline runs and logistics cost/profit history",
            "Scale path (optional): Kafka topics + Databricks/PySpark bronze→gold",
        ],
        0.6,
        1.5,
        12.1,
        5.0,
        size=15,
    )
    footer(s, 5, TOTAL)

    # 6. Data Source
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Data Source Breakdown", "Original data vs synthetic data vs live data")
    card(s, 0.5, 1.45, 4.0, 5.05)
    card(s, 4.7, 1.45, 4.0, 5.05)
    card(s, 8.9, 1.45, 4.0, 5.05)
    blocks = [
        ("Original / Trusted", [
            "Demand_prediction CSVs",
            "Commodities price datasets",
            "Logistics & cost CSVs",
            "Trade price tables",
            "Trained model bundles",
        ]),
        ("Live External", [
            "GDELT trade/commodity news",
            "Google News RSS (IN)",
            "Optional NewsAPI key",
            "Filtered to our commodities",
            "Demand countries only (demand box)",
        ]),
        ("Synthetic / Support", [
            "Synthetic news feature rows",
            "Used for training coverage",
            "Not used as final demo truth",
            "Helps edge-case robustness",
            "Controlled, not live market",
        ]),
    ]
    for i, (h, items) in enumerate(blocks):
        x = 0.7 + i * 4.2
        ht = s.shapes.add_textbox(Inches(x), Inches(1.65), Inches(3.6), Inches(0.4))
        run(ht.text_frame.paragraphs[0], h, size=16, bold=True, color=TEAL if i != 2 else GOLD)
        bx = s.shapes.add_textbox(Inches(x), Inches(2.2), Inches(3.6), Inches(4.0))
        tf = bx.text_frame
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run(p, f"• {it}", size=14, color=INK)
            p.space_after = Pt(9)
    footer(s, 6, TOTAL)

    # 7. Technical Workflow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Technical Workflow", "End-to-end multi-agent pipeline")
    bullets(
        s,
        [
            "1) News Agent fetches Demand news + India price news in parallel",
            "2) Demand Agent predicts top country–commodity demand scores (0–1)",
            "3) Price Agent forecasts next-month INR/quintal with news adjustment",
            "4) Logistics Agent optimizes India port → destination port corridors",
            "5) Container Agent allocates limited 20FT/40FT capacity by priority",
            "6) Explain Agent (Groq) + RAG Assistant produce reasoning and Q&A",
            "Persisted in SQLite → Dashboard cards + supervisor recommendation",
        ],
        0.6,
        1.5,
        12.1,
        5.0,
        size=16,
    )
    footer(s, 7, TOTAL)

    # 8. Architecture diagram
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "System Architecture & Data Flow", "Agents, APIs, models, and knowledge layer")
    if IMG_ARCH.exists():
        s.shapes.add_picture(str(IMG_ARCH), Inches(0.45), Inches(1.4), Inches(12.4), Inches(5.35))
    footer(s, 8, TOTAL)

    # 9. Existing vs Ours
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Existing System vs Our System", "Impact analysis & competitive edge")
    card(s, 0.5, 1.45, 6.0, 5.1)
    card(s, 6.8, 1.45, 6.0, 5.1)
    lh = s.shapes.add_textbox(Inches(0.75), Inches(1.65), Inches(5.5), Inches(0.4))
    run(lh.text_frame.paragraphs[0], "Existing / Typical Approach", size=16, bold=True, color=GOLD)
    rh = s.shapes.add_textbox(Inches(7.05), Inches(1.65), Inches(5.5), Inches(0.4))
    run(rh.text_frame.paragraphs[0], "ExportIntel AI", size=16, bold=True, color=TEAL)
    left = [
        "Spreadsheets + fragmented tools",
        "Manual news copy-paste",
        "Demand/price/logistics siloed",
        "Route profitability checked late",
        "Container plan is ad-hoc",
        "Hard to explain decisions",
    ]
    right = [
        "One-click multi-agent pipeline",
        "Auto live news (GDELT + RSS)",
        "Linked demand → price → profit",
        "Corridors ranked by INR net profit",
        "Priority-based container allocation",
        "LLM explanations + RAG Q&A",
    ]
    lb = s.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(5.5), Inches(4.0))
    rb = s.shapes.add_textbox(Inches(7.05), Inches(2.2), Inches(5.5), Inches(4.0))
    for box, items in ((lb, left), (rb, right)):
        tf = box.text_frame
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run(p, f"• {it}", size=15, color=INK)
            p.space_after = Pt(10)
    footer(s, 9, TOTAL)

    # 10. Container Prioritization (dedicated)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Container Prioritization", "How scarce capacity is allocated across export lanes")
    bullets(
        s,
        [
            "Problem: exporters have limited containers (e.g., 6 × 20FT) but multiple profitable opportunities.",
            "Container Agent ranks lanes using a blended priority score: demand + net profit + logistics cost + transit time.",
            "Example split: High priority lane gets more containers (e.g., 3), next gets 2, lower priority gets 1.",
            "Quantity is applied here: profit/container = profit/ton × payload tons; combined profit uses allocated containers.",
            "Output: export-first recommendation, allocation table, remaining capacity, and expected combined INR profit.",
            "UI module: Container Priority page shows ranked cards + full allocation plan after Run Analysis.",
        ],
        0.6,
        1.5,
        12.1,
        5.0,
        size=15,
    )
    footer(s, 10, TOTAL)

    # 11. Performance metrics
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "System Performance", "Key metrics & product benchmarks")
    bullets(
        s,
        [
            "End-to-end analysis: Demand + Price + Logistics + Container Prioritization in one Run Analysis.",
            "Live news fetch: parallel GDELT + Google RSS with relevance filters (max 10 headlines each).",
            "Price response to news: bullish vs bearish India news changes next-month prediction (not hard-coded).",
            "Profit engine: buy price always included; quantity applied via container payload × allocation.",
            "Persistence: SQLite stores full pipeline JSON + logistics cost/profit rows for audit.",
            "Local deploy: Frontend :5173 · Backend Flask :5001 (review-demo ready).",
        ],
        0.6,
        1.5,
        12.1,
        5.0,
        size=15,
    )
    footer(s, 11, TOTAL)

    # 12. Live Demo screens
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Live Demo Screens", "Working product samples from the deployed local app")
    imgs = [IMG3, IMG2, IMG1, IMG4]
    labels = ["Dashboard", "Demand Prediction", "Logistics Optimisation", "Container Priority"]
    coords = [(0.45, 1.45), (6.75, 1.45), (0.45, 4.15), (6.75, 4.15)]
    for img, (x, y), lab in zip(imgs, coords, labels):
        cap = s.shapes.add_textbox(Inches(x), Inches(y - 0.22), Inches(3.2), Inches(0.22))
        run(cap.text_frame.paragraphs[0], lab, size=11, bold=True, color=MUTED)
        if img.exists():
            s.shapes.add_picture(str(img), Inches(x), Inches(y), Inches(6.0), Inches(2.25))
    footer(s, 12, TOTAL)

    # 13. Alternative Approach
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Alternative Approach", "What we considered and why we chose this path")
    bullets(
        s,
        [
            "Alternative A — Pure LLM decisions without ML models: rejected (less stable, harder to audit).",
            "Alternative B — Manual pasted news only: rejected (not product-like for daily export ops).",
            "Alternative C — Kafka + Databricks first for every demo: deferred (excellent for scale, heavy for review demo).",
            "Selected approach — Hybrid ML + live news APIs + Flask/React UX, with Kafka/Databricks as roadmap.",
            "Why it wins — Reliable demos, explainable scores, INR profit clarity, clear production scale path.",
        ],
        0.6,
        1.5,
        12.1,
        5.0,
        size=16,
    )
    footer(s, 13, TOTAL)

    # 14. Future Scope
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Future Scope & Enhancements", "Roadmap for scalability and upgrades")
    bullets(
        s,
        [
            "Always-on Kafka news producers into Databricks bronze → silver → gold tables.",
            "Dynamic FX feeds and confidence intervals on demand/price predictions.",
            "What-if simulator for containers, port constraints, and tariff shocks.",
            "Multi-user roles, saved playbooks, and board-ready PDF/PPT exports.",
            "Backtesting harness on historical news weeks vs realized route outcomes.",
        ],
        0.6,
        1.5,
        12.1,
        5.0,
        size=16,
    )
    footer(s, 14, TOTAL)

    # 15. Impact / value (supports pitch style within structure)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Practical Impact", "Why this product matters in 14 minutes")
    bullets(
        s,
        [
            "Faster export decisions: one pipeline replaces fragmented analysis.",
            "Profit-first logistics: buy price + freight + route compared in INR.",
            "Capacity realism: scarce containers allocated by opportunity priority.",
            "Trust layer: Groq explanations + RAG grounded in project datasets.",
            "Demo outcome: live news → ranked opportunity → container plan → action.",
        ],
        0.6,
        1.55,
        12.1,
        4.9,
        size=17,
    )
    footer(s, 15, TOTAL)

    # 16. Speaking distribution
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    heading(s, "Team Participation Plan (~14 min)", "Equal speaking distribution for review")
    bullets(
        s,
        [
            "Member A (~3.5 min): Title/Team, Problem, Proposed Solution, Users",
            "Member B (~3.5 min): Tech Stack, Data Sources, Workflow, Architecture",
            "Member C (~3.5 min): Existing vs Ours, Container Prioritization, Live Demo",
            "Member D (~3.5 min): Alternative Approach, Future Scope, Conclusion + Q&A",
            "Demo tip: Run Analysis → Demand → Logistics → Container Priority → Recommendation.",
        ],
        0.6,
        1.55,
        12.1,
        4.9,
        size=16,
    )
    footer(s, 16, TOTAL)

    # 17. Conclusion & Q&A
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs, NAVY)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.15), prs.slide_width, Inches(3.2))
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL
    band.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.45), Inches(12), Inches(0.7))
    run(t.text_frame.paragraphs[0], "Conclusion & Q&A", size=34, bold=True, color=WHITE)
    st = s.shapes.add_textbox(Inches(0.7), Inches(3.25), Inches(12), Inches(1.4))
    run(
        st.text_frame.paragraphs[0],
        "ExportIntel AI turns live market news into ranked demand, price forecasts,\nprofitable routes, and prioritized container plans — with explainable AI support.",
        size=17,
        color=WHITE,
    )
    q = s.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(12), Inches(0.5))
    run(q.text_frame.paragraphs[0], "Thank you  ·  We are ready for questions", size=18, bold=True, color=RGBColor(220, 240, 245))
    footer(s, 17, TOTAL)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
