"""
ExportIntel AI Pitch PPT
- Slide 1: title only (reference split: ship + title text)
- Slide 2: team names / overview
- Other slides: standard layouts with SMALL images (not half-page)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "ExportIntel_AI_Final_Pitch.pptx"
ASSETS = ROOT / "assets" / "ppt"

TITLE_IMG = ASSETS / "title_ref.png"
SHIP = ASSETS / "ship_hero.png"
PORT = ASSETS / "port_terminal.png"
CONT = ASSETS / "containers_yard.png"
OCEAN = ASSETS / "ocean_trade.png"
ARCH = ROOT / "assets" / "ppt" / "architecture_chroma.png"
ARCH_FALLBACK = ROOT / "ExportIntel_AI_Architecture_Diagram.png"
DEMO_D = ASSETS / "demo_dashboard.png"
DEMO_DEM = ASSETS / "demo_demand.png"
DEMO_L = ASSETS / "demo_logistics.png"
DEMO_C = ASSETS / "demo_containers.png"

NAVY = RGBColor(15, 35, 70)
TEAL = RGBColor(20, 110, 130)
INK = RGBColor(30, 34, 42)
MUTED = RGBColor(95, 102, 112)
SOFT = RGBColor(248, 249, 251)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(220, 226, 234)
WARM = RGBColor(170, 105, 45)


def run(p, text, *, size=17, bold=False, color=INK, name="Calibri", italic=False):
    p.clear()
    r = p.add_run()
    r.text = text
    r.font.name = name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def paint(slide, prs, color=SOFT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def top_bar(slide, prs):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.16))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.16), prs.slide_width, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()


def footer(slide, n, total):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.1), Inches(12.2), Inches(0.25))
    run(box.text_frame.paragraphs[0], f"ExportIntel AI  ·  Review Pitch                              {n}/{total}", size=11, color=MUTED)


def heading(slide, title, subtitle=""):
    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
    run(t.text_frame.paragraphs[0], title, size=28, bold=True, color=NAVY, name="Calibri")
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(0.95), Inches(12.2), Inches(0.35))
        run(s.text_frame.paragraphs[0], subtitle, size=14, color=MUTED, italic=True)


def card(slide, x, y, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = LINE
    try:
        sh.adjustments[0] = 0.05
    except Exception:
        pass
    return sh


def bullets(slide, items, x, y, w, h, *, size=16):
    card(slide, x - 0.05, y - 0.08, w + 0.1, h + 0.16)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run(p, f"•  {item}", size=size, color=INK)
        p.space_after = Pt(8)


def small_image(slide, img_path, x, y, w=3.2, h=2.0):
    """Standard small image (not half-page)."""
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), Inches(w), Inches(h))
        # thin border frame look via shape behind? picture alone is fine
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        ph.fill.solid()
        ph.fill.fore_color.rgb = NAVY
        ph.line.fill.background()


def standard_content(prs, title, subtitle, items, n, total, *, img=None, size=16):
    """Text-first standard slide; optional small image on the right."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs)
    top_bar(s, prs)
    heading(s, title, subtitle)
    if img is not None and Path(img).exists():
        bullets(s, items, 0.6, 1.5, 8.7, 5.1, size=size)
        small_image(s, Path(img), 9.7, 1.7, w=3.1, h=2.1)
        # optional second small if tall enough - skip
    else:
        bullets(s, items, 0.6, 1.5, 12.1, 5.1, size=size)
    footer(s, n, total)
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    TOTAL = 18

    # ---------- 1. TITLE ONLY (reference style) ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs, WHITE)
    img = TITLE_IMG if TITLE_IMG.exists() else SHIP
    if img.exists():
        s.shapes.add_picture(str(img), Inches(0), Inches(0), Inches(6.7), Inches(7.5))
    else:
        ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(6.7), Inches(7.5))
        ph.fill.solid()
        ph.fill.fore_color.rgb = NAVY
        ph.line.fill.background()

    # right panel title only — centered like reference
    title = s.shapes.add_textbox(Inches(7.0), Inches(2.55), Inches(5.8), Inches(2.2))
    tf = title.text_frame
    tf.word_wrap = True
    run(
        tf.paragraphs[0],
        "Predictive Export Intelligence\n& Logistics Optimization\nPlatform",
        size=28,
        bold=True,
        color=RGBColor(20, 20, 20),
        name="Calibri",
    )
    # tiny brand line
    brand = s.shapes.add_textbox(Inches(7.0), Inches(5.1), Inches(5.8), Inches(0.35))
    run(brand.text_frame.paragraphs[0], "ExportIntel AI", size=14, color=TEAL, bold=True)

    # ---------- 2. TEAM OVERVIEW (names here) ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs)
    top_bar(s, prs)
    heading(s, "Title & Team Overview", "Team details for review submission")
    card(s, 0.55, 1.5, 8.8, 5.1)
    info = s.shapes.add_textbox(Inches(0.85), Inches(1.75), Inches(8.2), Inches(4.6))
    tf = info.text_frame
    run(tf.paragraphs[0], "Project: Predictive Export Intelligence & Logistics Optimization Platform", size=15, bold=True, color=NAVY)
    for line in [
        "",
        "Team Number: __________",
        "Team Name: ________________________________",
        "Project Leader: ____________________________",
        "",
        "Team Members:",
        "1) ________________________________________",
        "2) ________________________________________",
        "3) ________________________________________",
        "4) ________________________________________",
        "",
        "Tech Stack (form): React, Flask, XGBoost, Groq LLM, Chroma, GDELT/Google News, SQLite",
    ]:
        p = tf.add_paragraph()
        run(p, line, size=15, color=INK)
        p.space_after = Pt(3)
    small_image(s, CONT if CONT.exists() else PORT, 9.7, 1.7, 3.1, 2.1)
    small_image(s, PORT if PORT.exists() else OCEAN, 9.7, 4.1, 3.1, 2.1)
    footer(s, 2, TOTAL)

    # ---------- 3. Problem ----------
    standard_content(
        prs,
        "Problem Statement",
        "Core problem the project solves",
        [
            "Indian exporters must decide what to export, where to sell, and which route is profitable under fast market shifts.",
            "Demand signals, mandi prices, freight costs, and container limits are scattered across tools.",
            "Manual planning is slow and often misses live news that changes opportunity overnight.",
            "Without a unified system, teams risk wrong country–commodity choices and loss-making lanes.",
        ],
        3,
        TOTAL,
        img=PORT,
        size=16,
    )

    # ---------- 4. Solution ----------
    standard_content(
        prs,
        "Proposed Solution",
        "High-level overview of ExportIntel AI",
        [
            "A multi-agent React + Flask platform for end-to-end export decisions.",
            "One click: Live News → Demand → Price → Logistics → Container Prioritization.",
            "ML models score demand opportunities and forecast next-month prices (INR/quintal).",
            "Logistics ranks India→destination corridors by net profit (INR @ ₹96.3/USD).",
            "Groq LLM explanations + RAG Trade Assistant make decisions explainable.",
        ],
        4,
        TOTAL,
        img=OCEAN,
        size=15,
    )

    # ---------- 5. Users ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs)
    top_bar(s, prs)
    heading(s, "Target Users & Use Cases", "Intended audience and practical scenarios")
    cols = [
        ("Export Analysts", ["Rank demand opportunities", "Review next-month buy prices", "Compare profitable corridors"]),
        ("Operations Teams", ["Allocate limited containers", "Prioritize export-first lanes", "Track cost, transit, profit"]),
        ("Decision Makers", ["Use supervisor recommendation", "Validate INR profit impact", "Ask Trade Assistant (RAG)"]),
    ]
    for i, (h, items) in enumerate(cols):
        x = 0.5 + i * 4.2
        card(s, x, 1.5, 4.0, 5.0)
        ht = s.shapes.add_textbox(Inches(x + 0.2), Inches(1.75), Inches(3.6), Inches(0.45))
        run(ht.text_frame.paragraphs[0], h, size=17, bold=True, color=TEAL)
        bx = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.4), Inches(3.6), Inches(3.8))
        tf = bx.text_frame
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run(p, f"• {it}", size=15, color=INK)
            p.space_after = Pt(12)
    footer(s, 5, TOTAL)

    # ---------- 6. Tech ----------
    standard_content(
        prs,
        "Tech Stack",
        "Technologies used in the platform",
        [
            "Frontend: React + Vite  |  Backend: Flask REST APIs",
            "ML: XGBoost / Joblib + Pandas pipelines",
            "LLM: Groq Llama 3.3 70B (explanations + RAG answers)",
            "Vector DB: Chroma (embeddings for Trade Assistant RAG)",
            "Live News: GDELT + Google News RSS (+ optional NewsAPI)",
            "Persistence: SQLite for pipeline and logistics history",
            "Scale path: Kafka + Databricks / PySpark (roadmap)",
        ],
        6,
        TOTAL,
        img=SHIP,
        size=15,
    )

    # ---------- 7. Data ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs)
    top_bar(s, prs)
    heading(s, "Data Source Breakdown", "Original vs Live vs Synthetic")
    blocks = [
        ("Original / Trusted", TEAL, ["Demand_prediction CSVs", "Commodities price data", "Logistics & cost CSVs", "Trade price tables", "Trained model bundles"]),
        ("Live External", NAVY, ["GDELT commodity news", "Google News RSS (IN)", "Optional NewsAPI", "Filtered commodities", "Demand countries only"]),
        ("Synthetic Support", WARM, ["Synthetic feature rows", "Training coverage aid", "Not final demo truth", "Edge-case robustness", "Controlled datasets"]),
    ]
    for i, (h, col, items) in enumerate(blocks):
        x = 0.5 + i * 4.2
        card(s, x, 1.5, 4.0, 5.0)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(4.0), Inches(0.1))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = col
        stripe.line.fill.background()
        ht = s.shapes.add_textbox(Inches(x + 0.2), Inches(1.8), Inches(3.6), Inches(0.4))
        run(ht.text_frame.paragraphs[0], h, size=16, bold=True, color=col)
        bx = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.4), Inches(3.6), Inches(3.8))
        tf = bx.text_frame
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run(p, f"• {it}", size=14, color=INK)
            p.space_after = Pt(10)
    footer(s, 7, TOTAL)

    # ---------- 8. Workflow ----------
    standard_content(
        prs,
        "Technical Workflow",
        "System pipeline and data flow",
        [
            "1) News Agent — demand news + India price news",
            "2) Demand Agent — country–commodity demand scores",
            "3) Price Agent — next-month INR/quintal (news-adjusted)",
            "4) Logistics Agent — port corridor + net profit",
            "5) Container Prioritization — allocate scarce capacity",
            "6) Explain Agent + RAG — reasoning and Q&A",
            "Persisted in SQLite → dashboard recommendation",
        ],
        8,
        TOTAL,
        img=CONT,
        size=15,
    )

    # ---------- 9. Architecture ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs)
    top_bar(s, prs)
    heading(s, "System Architecture", "Architecture diagram and agent flow (RAG via Chroma DB)")
    arch_img = ARCH if ARCH.exists() else ARCH_FALLBACK
    if arch_img.exists():
        s.shapes.add_picture(str(arch_img), Inches(0.7), Inches(1.45), Inches(11.9), Inches(5.2))
    footer(s, 9, TOTAL)

    # ---------- 10. Five core modules ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs)
    top_bar(s, prs)
    heading(s, "Five Core Modules", "Demand · Commodity Price · Logistics · Container · RAG")
    modules = [
        ("Demand", "Scores country–commodity export opportunities from live news + ML demand model."),
        ("Commodity Price", "Forecasts next-month India buy price (INR/quintal) with news sentiment adjustment."),
        ("Logistics Optimization", "Ranks India→destination corridors by net profit in INR (FX @ ₹96.3/USD)."),
        ("Container Prioritization", "Allocates limited containers across lanes by demand, profit, cost, and transit."),
        ("RAG (Chroma DB)", "Trade Assistant retrieves trusted chunks from Chroma vector DB; Groq answers from those chunks only."),
    ]
    for i, (h, text) in enumerate(modules):
        y = 1.45 + i * 1.05
        card(s, 0.55, y, 12.2, 0.95)
        ht = s.shapes.add_textbox(Inches(0.8), Inches(y + 0.12), Inches(11.7), Inches(0.32))
        run(ht.text_frame.paragraphs[0], f"{i + 1}. {h}", size=16, bold=True, color=TEAL)
        bx = s.shapes.add_textbox(Inches(0.8), Inches(y + 0.45), Inches(11.7), Inches(0.4))
        run(bx.text_frame.paragraphs[0], text, size=14, color=INK)
    footer(s, 10, TOTAL)

    # ---------- 11. Existing vs Ours ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs)
    top_bar(s, prs)
    heading(s, "Existing System vs Our System", "Impact analysis and competitive edge")
    card(s, 0.5, 1.45, 6.0, 5.15)
    card(s, 6.8, 1.45, 6.0, 5.15)
    lh = s.shapes.add_textbox(Inches(0.75), Inches(1.65), Inches(5.5), Inches(0.4))
    run(lh.text_frame.paragraphs[0], "Existing / Typical Approach", size=16, bold=True, color=WARM)
    rh = s.shapes.add_textbox(Inches(7.05), Inches(1.65), Inches(5.5), Inches(0.4))
    run(rh.text_frame.paragraphs[0], "ExportIntel AI", size=16, bold=True, color=TEAL)
    left = ["Spreadsheets + fragmented tools", "Manual news copy-paste", "Demand/price/logistics siloed", "Route profitability checked late", "Container plan is ad-hoc", "Hard to explain decisions"]
    right = ["One-click multi-agent pipeline", "Auto live news (GDELT + RSS)", "Linked demand → price → profit", "Corridors ranked by INR net profit", "Priority-based container allocation", "LLM explanations + RAG Q&A"]
    lb = s.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(5.5), Inches(4.1))
    rb = s.shapes.add_textbox(Inches(7.05), Inches(2.2), Inches(5.5), Inches(4.1))
    for box, items in ((lb, left), (rb, right)):
        tf = box.text_frame
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run(p, f"• {it}", size=15, color=INK)
            p.space_after = Pt(10)
    footer(s, 11, TOTAL)

    # ---------- 12. Container Prioritization (important) ----------
    standard_content(
        prs,
        "Container Prioritization",
        "How limited containers are allocated across export lanes",
        [
            "Exporters often have limited containers (e.g., 6 × 20FT) but multiple opportunities.",
            "Container Agent ranks lanes using demand + net profit + logistics cost + transit.",
            "Example allocation: high priority 3, medium 2, lower 1 (until capacity is used).",
            "Quantity is applied here: profit/container = profit/ton × payload tons.",
            "Output: export-first recommendation, allocation table, combined INR profit.",
            "Shown on Container Priority page after Run Analysis.",
        ],
        12,
        TOTAL,
        img=CONT,
        size=15,
    )

    # ---------- 13. Performance ----------
    standard_content(
        prs,
        "System Performance & Live Demo",
        "Metrics, benchmarks, and demo readiness",
        [
            "Full pipeline in one Run Analysis: Demand + Price + Logistics + Containers.",
            "Live news: parallel GDELT + Google RSS, filtered, max 10 headlines each.",
            "Price changes with bullish/bearish India news (not hard-coded).",
            "Buy price included in profit; quantity used at container allocation.",
            "SQLite stores pipeline JSON + logistics profit history.",
            "Local demo: Frontend :5173 · Backend Flask :5001.",
        ],
        13,
        TOTAL,
        img=PORT,
        size=15,
    )

    # ---------- 14. Demo screens ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs)
    top_bar(s, prs)
    heading(s, "Live Product Screens", "Dashboard · Demand · Logistics · Container Priority")
    imgs = [DEMO_D, DEMO_DEM, DEMO_L, DEMO_C]
    labels = ["Dashboard", "Demand", "Logistics", "Container Priority"]
    coords = [(0.5, 1.5), (6.8, 1.5), (0.5, 4.2), (6.8, 4.2)]
    for img, (x, y), lab in zip(imgs, coords, labels):
        cap = s.shapes.add_textbox(Inches(x), Inches(y - 0.22), Inches(3.2), Inches(0.2))
        run(cap.text_frame.paragraphs[0], lab, size=12, bold=True, color=MUTED)
        if img.exists():
            s.shapes.add_picture(str(img), Inches(x), Inches(y), Inches(5.9), Inches(2.2))
    footer(s, 14, TOTAL)

    # ---------- 15. Alternative ----------
    standard_content(
        prs,
        "Alternative Approach",
        "Secondary options considered and why primary was selected",
        [
            "Pure LLM decisions without ML — rejected (unstable, hard to audit).",
            "Manual pasted news only — rejected (not product-like for daily ops).",
            "Kafka + Databricks first for every demo — deferred (scale later).",
            "Selected: Hybrid ML + live news APIs + Flask/React UX.",
            "Kafka/Databricks kept as production streaming roadmap.",
        ],
        15,
        TOTAL,
        img=OCEAN,
        size=15,
    )

    # ---------- 16. Future ----------
    standard_content(
        prs,
        "Future Scope & Enhancements",
        "Roadmap for scalability and upgrades",
        [
            "Kafka producers → Databricks bronze / silver / gold.",
            "Dynamic FX and prediction confidence intervals.",
            "What-if simulator for containers, ports, and tariffs.",
            "Multi-user roles and board-ready export packs.",
            "Backtesting on historical news weeks vs outcomes.",
        ],
        16,
        TOTAL,
        img=PORT,
        size=16,
    )

    # ---------- 17. Speaking plan ----------
    standard_content(
        prs,
        "Team Participation Plan (~14 min)",
        "Equal speaking distribution",
        [
            "Member A (~3.5 min): Title, Team, Problem, Solution, Users",
            "Member B (~3.5 min): Tech Stack, Data, Workflow, Architecture",
            "Member C (~3.5 min): Comparison, Container Prioritization, Live Demo",
            "Member D (~3.5 min): Alternatives, Future Scope, Conclusion + Q&A",
            "Demo tip: Run Analysis → Demand → Logistics → Container Priority.",
        ],
        17,
        TOTAL,
        size=16,
    )

    # ---------- 18. Conclusion ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint(s, prs, NAVY)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), prs.slide_width, Inches(3.2))
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL
    band.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(0.7))
    run(t.text_frame.paragraphs[0], "Conclusion & Q&A", size=34, bold=True, color=WHITE)
    st = s.shapes.add_textbox(Inches(0.7), Inches(3.35), Inches(12), Inches(1.2))
    run(
        st.text_frame.paragraphs[0],
        "ExportIntel AI turns live market news into ranked demand, price forecasts,\nprofitable routes, and prioritized container plans — with explainable AI.",
        size=17,
        color=WHITE,
    )
    q = s.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(12), Inches(0.4))
    run(q.text_frame.paragraphs[0], "Thank you  ·  Ready for questions", size=18, bold=True, color=RGBColor(210, 235, 245))
    # small image accent bottom-right
    if SHIP.exists():
        small_image(s, SHIP, 10.3, 5.7, 2.5, 1.2)
    footer(s, 18, TOTAL)

    out = OUT
    try:
        prs.save(str(out))
    except PermissionError:
        out = ROOT / "ExportIntel_AI_Final_Pitch_v2.pptx"
        prs.save(str(out))
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
