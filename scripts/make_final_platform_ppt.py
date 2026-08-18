"""
ExportIntel AI — Final Presentation
Visual language of ExportIntel_AI_Platform.pptx + full content from Presentation_v2.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "ExportIntel_AI_Final.pptx"
OUT_DL = Path(r"c:\Users\Lenovo\Downloads\ExportIntel_AI_Final.pptx")

ASSETS = ROOT / "assets" / "ppt"
PLAT = ASSETS / "platform_slides"
ARCH = ASSETS / "architecture_chroma.png"
ARCH_V2 = ASSETS / "v2_slide_05.png"
ARCH_FALLBACK = ROOT / "ExportIntel_AI_Architecture_Diagram.png"
DEMO_D = ASSETS / "demo_dashboard.png"
DEMO_DEM = ASSETS / "demo_demand.png"
DEMO_L = ASSETS / "demo_logistics.png"
DEMO_C = ASSETS / "demo_containers.png"
SHIP = ASSETS / "ship_hero.png"
PORT = ASSETS / "port_terminal.png"
CONT = ASSETS / "containers_yard.png"
OCEAN = ASSETS / "ocean_trade.png"

# Platform-inspired palette
NAVY = RGBColor(8, 18, 40)
NAVY2 = RGBColor(12, 28, 58)
PANEL = RGBColor(16, 36, 68)
INK = RGBColor(230, 236, 245)
MUTED = RGBColor(160, 175, 195)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(30, 190, 200)
GOLD = RGBColor(230, 170, 55)
CORAL = RGBColor(230, 90, 70)
SOFT_LINE = RGBColor(50, 80, 120)

SLIDE_W = 13.333
SLIDE_H = 7.5
TOTAL = 21


def run(p, text, *, size=16, bold=False, color=INK, name="Calibri", italic=False, align=None):
    p.clear()
    r = p.add_run()
    r.text = text
    r.font.name = name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if align is not None:
        p.alignment = align


def dark_bg(slide, prs, photo=None, dim=0.78):
    """Dark navy base with optional dimmed photo atmosphere."""
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    base.fill.solid()
    base.fill.fore_color.rgb = NAVY
    base.line.fill.background()
    if photo and Path(photo).exists():
        slide.shapes.add_picture(str(photo), 0, 0, prs.slide_width, prs.slide_height)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = NAVY
        overlay.line.fill.background()
        try:
            from lxml import etree
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            solidFill = overlay._element.spPr.find(f"{ns}solidFill")
            srgb = solidFill.find(f"{ns}srgbClr") if solidFill is not None else None
            if srgb is not None:
                alpha = etree.SubElement(srgb, f"{ns}alpha")
                alpha.set("val", str(int(dim * 100000)))
        except Exception:
            pass
    # thin teal top accent (Platform-like frame)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL
    top.line.fill.background()


def full_image_slide(prs, img_path):
    """Drop a Platform-rendered slide in as a full-bleed image."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs)
    if Path(img_path).exists():
        s.shapes.add_picture(str(img_path), 0, 0, prs.slide_width, prs.slide_height)
    return s


def glass(slide, x, y, w, h, *, fill=PANEL, line=SOFT_LINE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1.25)
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh


def accent_bar(slide, x, y, h, color=TEAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def footer(slide, n):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.28))
    run(
        box.text_frame.paragraphs[0],
        f"ExportIntel AI  ·  Predictive Export Intelligence                              {n} / {TOTAL}",
        size=11,
        color=MUTED,
    )


def heading(slide, title, subtitle=""):
    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.55))
    run(t.text_frame.paragraphs[0], title, size=28, bold=True, color=WHITE)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(0.88), Inches(12.2), Inches(0.35))
        run(s.text_frame.paragraphs[0], subtitle, size=14, color=TEAL, italic=True)


def bullets_in(slide, items, x, y, w, h, *, size=15, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run(p, f"•  {item}", size=size, color=color)
        p.space_after = Pt(10)
    return box


def content_slide(prs, title, subtitle, items, n, *, photo=None, size=15, accent=TEAL):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, photo, dim=0.72)
    heading(s, title, subtitle)
    glass(s, 0.5, 1.4, 12.3, 5.3)
    accent_bar(s, 0.5, 1.4, 5.3, accent)
    bullets_in(s, items, 0.85, 1.65, 11.7, 4.8, size=size)
    footer(s, n)
    return s


def two_col_slide(prs, title, subtitle, left_title, left_items, right_title, right_items, n, *, photo=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, photo, dim=0.72)
    heading(s, title, subtitle)
    glass(s, 0.5, 1.4, 6.0, 5.3)
    accent_bar(s, 0.5, 1.4, 5.3, TEAL)
    ht = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.4), Inches(0.4))
    run(ht.text_frame.paragraphs[0], left_title, size=16, bold=True, color=TEAL)
    bullets_in(s, left_items, 0.8, 2.15, 5.4, 4.2, size=14)

    glass(s, 6.8, 1.4, 6.0, 5.3)
    accent_bar(s, 6.8, 1.4, 5.3, GOLD)
    ht2 = s.shapes.add_textbox(Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.4))
    run(ht2.text_frame.paragraphs[0], right_title, size=16, bold=True, color=GOLD)
    bullets_in(s, right_items, 7.1, 2.15, 5.4, 4.2, size=14)
    footer(s, n)
    return s


def cards_row(prs, title, subtitle, cards, n, *, photo=None):
    """cards: list of (title, items, accent_color)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, photo, dim=0.72)
    heading(s, title, subtitle)
    n_cards = len(cards)
    gap = 0.25
    margin = 0.5
    usable = SLIDE_W - 2 * margin - gap * (n_cards - 1)
    w = usable / n_cards
    for i, (h, items, col) in enumerate(cards):
        x = margin + i * (w + gap)
        glass(s, x, 1.45, w, 5.2)
        accent_bar(s, x, 1.45, 5.2, col)
        ht = s.shapes.add_textbox(Inches(x + 0.25), Inches(1.7), Inches(w - 0.4), Inches(0.45))
        run(ht.text_frame.paragraphs[0], h, size=15, bold=True, color=col)
        bullets_in(s, items, x + 0.25, 2.3, w - 0.45, 4.0, size=13)
    footer(s, n)
    return s


def four_phase_slide(prs, n):
    """Chaos → Clarity with the four decision phases."""
    phases = [
        ("01", "Commodity Price\nPrediction", "Next-month India buy\nprice (INR/quintal)", GOLD),
        ("02", "Demand\nPrediction", "Country–commodity\nopportunity scores", TEAL),
        ("03", "Logistic\nOptimization", "Corridor routes &\nINR net profit", CORAL),
        ("04", "Container\nPriority", "Scarce capacity\nallocated by priority", TEAL),
    ]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, PORT, dim=0.78)
    heading(s, "Four-Phase Decision Pipeline", "Predictive Export Intelligence — from market signal to export action")

    # left chaos / right clarity labels
    left = s.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(2.2), Inches(0.35))
    run(left.text_frame.paragraphs[0], "CHAOS  →", size=13, bold=True, color=MUTED)
    right = s.shapes.add_textbox(Inches(10.6), Inches(1.35), Inches(2.2), Inches(0.35))
    run(right.text_frame.paragraphs[0], "→  CLARITY", size=13, bold=True, color=TEAL)
    right.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    gap = 0.28
    margin = 0.5
    usable = SLIDE_W - 2 * margin - gap * 3
    w = usable / 4
    y, h = 1.85, 4.55
    for i, (num, title, blurb, col) in enumerate(phases):
        x = margin + i * (w + gap)
        glass(s, x, y, w, h, fill=PANEL, line=col)
        accent_bar(s, x, y, h, col)

        nb = s.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.25), Inches(w - 0.4), Inches(0.45))
        run(nb.text_frame.paragraphs[0], f"PHASE {num}", size=12, bold=True, color=col)

        tb = s.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.85), Inches(w - 0.4), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for line in title.split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run(p, line, size=20, bold=True, color=WHITE)

        db = s.shapes.add_textbox(Inches(x + 0.22), Inches(y + 2.7), Inches(w - 0.4), Inches(1.4))
        dtf = db.text_frame
        dtf.word_wrap = True
        first = True
        for line in blurb.split("\n"):
            p = dtf.paragraphs[0] if first else dtf.add_paragraph()
            first = False
            run(p, line, size=13, color=MUTED)

        if i < 3:
            # small connector chevron between cards
            cx = x + w + 0.02
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(cx), Inches(y + h / 2 - 0.12), Inches(0.24), Inches(0.24)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = SOFT_LINE
            arr.line.fill.background()

    footer(s, n)
    return s


def title_slide(prs):
    """Predictive Export Intelligence title with four-phase strip."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    photo = SHIP if SHIP.exists() else (PLAT / "slide_01.png")
    dark_bg(s, prs, photo if Path(photo).exists() else None, dim=0.62)

    panel = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.35), Inches(8.6), Inches(3.55)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY2
    panel.line.color.rgb = TEAL
    try:
        panel.adjustments[0] = 0.04
    except Exception:
        pass
    accent_bar(s, 0.7, 1.35, 3.55, CORAL)

    brand = s.shapes.add_textbox(Inches(1.05), Inches(1.6), Inches(7.9), Inches(0.4))
    run(brand.text_frame.paragraphs[0], "EXPORTINTEL AI", size=14, bold=True, color=TEAL)

    title = s.shapes.add_textbox(Inches(1.05), Inches(2.1), Inches(7.9), Inches(1.3))
    tf = title.text_frame
    tf.word_wrap = True
    run(tf.paragraphs[0], "Predictive Export Intelligence", size=32, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    run(p2, "& Logistics Optimization", size=26, bold=True, color=WHITE)

    sub = s.shapes.add_textbox(Inches(1.05), Inches(3.7), Inches(7.9), Inches(0.8))
    run(
        sub.text_frame.paragraphs[0],
        "AI-powered decision support for Indian commodity exporters —\nfrom live market signals to profitable, prioritized shipments.",
        size=14,
        color=MUTED,
    )

    # four-phase strip at bottom
    labels = [
        ("01", "Price Prediction", GOLD),
        ("02", "Demand Prediction", TEAL),
        ("03", "Logistic Optimization", CORAL),
        ("04", "Container Priority", TEAL),
    ]
    gap, margin, usable = 0.22, 0.7, SLIDE_W - 1.4 - 0.66
    w = usable / 4
    for i, (num, lab, col) in enumerate(labels):
        x = margin + i * (w + gap)
        glass(s, x, 5.35, w, 1.35, fill=PANEL, line=col)
        nb = s.shapes.add_textbox(Inches(x + 0.18), Inches(5.5), Inches(w - 0.3), Inches(0.35))
        run(nb.text_frame.paragraphs[0], f"PHASE {num}", size=11, bold=True, color=col)
        lb = s.shapes.add_textbox(Inches(x + 0.18), Inches(5.95), Inches(w - 0.3), Inches(0.55))
        run(lb.text_frame.paragraphs[0], lab, size=14, bold=True, color=WHITE)
    return s


def agent_pipeline_slide(prs, n):
    """Four core phases (+ news intake & explain support)."""
    phases = [
        ("News Intake", "GDELT + Google RSS\nlive market headlines", MUTED),
        ("1. Price", "Next-month India buy\nprice (INR/quintal)", GOLD),
        ("2. Demand", "Country–commodity\nopportunity scores", TEAL),
        ("3. Logistics", "Corridor routes &\nINR net profit", CORAL),
        ("4. Containers", "Priority allocation\nof scarce capacity", TEAL),
        ("Explain", "Groq LLM + RAG\nreasoning & Q&A", MUTED),
    ]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, CONT, dim=0.78)
    heading(
        s,
        "Predictive Export Intelligence Pipeline",
        "Four decision phases powered by live news and explainable AI",
    )
    gap, margin = 0.18, 0.4
    usable = SLIDE_W - 2 * margin - gap * 5
    w = usable / 6
    y, h = 1.55, 4.85
    for i, (title, blurb, col) in enumerate(phases):
        x = margin + i * (w + gap)
        glass(s, x, y, w, h)
        accent_bar(s, x, y, h, col)
        ht = s.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.35), Inches(w - 0.28), Inches(1.0))
        htf = ht.text_frame
        htf.word_wrap = True
        run(htf.paragraphs[0], title, size=15, bold=True, color=col)
        db = s.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.6), Inches(w - 0.28), Inches(2.6))
        dtf = db.text_frame
        dtf.word_wrap = True
        first = True
        for line in blurb.split("\n"):
            p = dtf.paragraphs[0] if first else dtf.add_paragraph()
            first = False
            run(p, line, size=12, color=INK)
        if i < 5:
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + w + 0.01),
                Inches(y + h / 2 - 0.1),
                Inches(0.16),
                Inches(0.2),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = SOFT_LINE
            arr.line.fill.background()
    footer(s, n)
    return s


def comparison_slide(prs, n):
    rows = [
        ("Market data", "Scattered tools & manual news copy-paste", "Live news feeds one predictive pipeline"),
        ("Price & demand", "Analyzed in silos, often out of sync", "Phase 1 price → Phase 2 demand, linked"),
        ("Logistics", "Route profit checked late", "Phase 3 ranks corridors by INR net profit"),
        ("Containers", "Ad-hoc allocation of scarce capacity", "Phase 4 priority-based export-first plan"),
        ("Explainability", "Hard to justify gut decisions", "Groq LLM explanations + RAG Q&A"),
    ]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs, PORT, dim=0.78)
    heading(s, "Existing Approach vs ExportIntel AI", "Why predictive export intelligence wins")

    # headers
    glass(s, 3.3, 1.35, 4.5, 0.55)
    glass(s, 8.0, 1.35, 4.8, 0.55, line=TEAL)
    h1 = s.shapes.add_textbox(Inches(3.45), Inches(1.45), Inches(4.2), Inches(0.4))
    run(h1.text_frame.paragraphs[0], "Spreadsheet Status Quo", size=14, bold=True, color=MUTED)
    h2 = s.shapes.add_textbox(Inches(8.15), Inches(1.45), Inches(4.5), Inches(0.4))
    run(h2.text_frame.paragraphs[0], "ExportIntel AI Platform", size=14, bold=True, color=TEAL)

    row_h = 0.88
    for i, (label, left, right) in enumerate(rows):
        y = 2.1 + i * row_h
        glass(s, 0.5, y, 2.6, row_h - 0.1)
        glass(s, 3.3, y, 4.5, row_h - 0.1)
        glass(s, 8.0, y, 4.8, row_h - 0.1, line=TEAL if i % 2 == 0 else CORAL)
        lb = s.shapes.add_textbox(Inches(0.65), Inches(y + 0.22), Inches(2.3), Inches(0.5))
        run(lb.text_frame.paragraphs[0], label.upper(), size=12, bold=True, color=GOLD)
        lt = s.shapes.add_textbox(Inches(3.45), Inches(y + 0.18), Inches(4.2), Inches(0.55))
        ltf = lt.text_frame
        ltf.word_wrap = True
        run(ltf.paragraphs[0], left, size=12, color=INK)
        rt = s.shapes.add_textbox(Inches(8.15), Inches(y + 0.18), Inches(4.5), Inches(0.55))
        rtf = rt.text_frame
        rtf.word_wrap = True
        run(rtf.paragraphs[0], right, size=12, color=WHITE)
    footer(s, n)
    return s


def thank_you_slide(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    photo = PLAT / "slide_01.png"
    if photo.exists():
        s.shapes.add_picture(str(photo), 0, 0, prs.slide_width, prs.slide_height)
        panel = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(1.85), Inches(9.7), Inches(3.7)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = NAVY2
        panel.line.color.rgb = TEAL
        try:
            panel.adjustments[0] = 0.05
        except Exception:
            pass
    else:
        dark_bg(s, prs)
    t = s.shapes.add_textbox(Inches(2.1), Inches(2.15), Inches(9.1), Inches(0.7))
    run(t.text_frame.paragraphs[0], "Thank You", size=42, bold=True, color=WHITE)
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    st = s.shapes.add_textbox(Inches(2.1), Inches(2.95), Inches(9.1), Inches(0.45))
    run(st.text_frame.paragraphs[0], "Predictive Export Intelligence", size=20, color=TEAL)
    st.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sub = s.shapes.add_textbox(Inches(2.1), Inches(3.55), Inches(9.1), Inches(1.2))
    run(
        sub.text_frame.paragraphs[0],
        "1 Price  →  2 Demand  →  3 Logistics  →  4 Containers\nQuestions and discussion",
        size=15,
        color=MUTED,
    )
    sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer(s, n)
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 1 — Title
    title_slide(prs)

    # 2 — Agenda
    content_slide(
        prs,
        "Agenda",
        "Predictive Export Intelligence — what we will cover",
        [
            "1. Problem and project objectives",
            "2. Four-phase pipeline & technology stack",
            "3. Price → Demand → Logistics → Container logic",
            "4. Live demo outcomes and business impact",
            "5. Results, value, and future roadmap",
        ],
        2,
        photo=PORT,
        size=18,
    )

    # 3 — Problem (Platform visual)
    full_image_slide(prs, PLAT / "slide_02.png")

    # 4 — Objectives
    content_slide(
        prs,
        "Project Objectives",
        "ExportIntel AI — Predictive Export Intelligence",
        [
            "1. Commodity price prediction — next-month India buy prices (INR/quintal)",
            "2. Demand prediction — top country–commodity opportunities from live signals",
            "3. Logistic optimization — India port → destination corridor net profit (INR)",
            "4. Container priority — allocate limited capacity across competing lanes",
            "Plus explainable insights and document-grounded Q&A (RAG + Groq LLM)",
        ],
        4,
        photo=OCEAN,
        size=16,
        accent=GOLD,
    )

    # 5 — Four-phase pipeline
    four_phase_slide(prs, 5)

    # 6 — Personas (Platform)
    full_image_slide(prs, PLAT / "slide_04.png")

    # 7 — Tech stack
    two_col_slide(
        prs,
        "Technology Stack",
        "Application layer + predictive intelligence layer",
        "Application Layer",
        [
            "Frontend: React + Vite dashboard",
            "Backend: Flask REST APIs (:5001)",
            "Data processing: Pandas",
            "Persistence: SQLite (pipeline + logistics history)",
            "Modules: Price · Demand · Logistics · Containers · RAG",
        ],
        "AI & Intelligence Layer",
        [
            "ML: XGBoost / Joblib price & demand models",
            "LLM: Groq Llama 3.3 70B",
            "RAG: Chroma Vector DB (+ TF-IDF fallback)",
            "News APIs: GDELT + Google News RSS",
            "Explain Agent for reasoning narratives",
        ],
        7,
        photo=CONT,
    )

    # 8 — Data sources (Platform)
    full_image_slide(prs, PLAT / "slide_05.png")

    # 9 — Agent / phase pipeline
    agent_pipeline_slide(prs, 9)

    # 10 — Software architecture (Platform)
    full_image_slide(prs, PLAT / "slide_07.png")

    # 11 — System architecture diagram
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s, prs)
    heading(s, "System Architecture Diagram", "Predictive Export Intelligence — end-to-end data flow")
    arch = ARCH if ARCH.exists() else (ARCH_V2 if ARCH_V2.exists() else ARCH_FALLBACK)
    if arch.exists():
        s.shapes.add_picture(str(arch), Inches(0.55), Inches(1.35), Inches(12.2), Inches(5.4))
    footer(s, 11)

    # 12 — Phases 1–2
    two_col_slide(
        prs,
        "Phases 1 & 2 — Price and Demand",
        "Predictive core of ExportIntel AI",
        "1. Commodity Price Prediction",
        [
            "Baseline from commodity mandi dataset (INR/quintal)",
            "XGBoost predicts base next-month price per commodity",
            "News sentiment applies controlled +/- adjustment",
            "Predictions can change across runs when news tone shifts",
        ],
        "2. Demand Prediction",
        [
            "News headlines parsed for country, commodity, and sentiment",
            "Model computes demand score (0–1) per country–commodity pair",
            "Shortage / production drop / export opportunity raise the score",
            "Top opportunities ranked and passed to logistics & containers",
        ],
        12,
        photo=PORT,
    )

    # 13 — Phases 3–4
    two_col_slide(
        prs,
        "Phases 3 & 4 — Logistics and Containers",
        "From predicted opportunity to executable export plan",
        "3. Logistic Optimization",
        [
            "Net profit/ton = sell − buy − logistics cost/ton",
            "Buy amount from predicted India commodity price",
            "Corridors ranked by INR net profit (FX: 1 USD = ₹96.3)",
            "RAG + live news (GDELT/RSS) ground explanations",
        ],
        "4. Container Priority",
        [
            "Limited 20FT/40FT capacity across competing lanes",
            "Priority from demand + profit + cost + transit",
            "Quantity applied: profit/container = profit/ton × payload",
            "Export-first allocation plan for operations teams",
        ],
        13,
        photo=OCEAN,
    )

    # 14 — Container prioritization (Platform)
    full_image_slide(prs, PLAT / "slide_08.png")

    # 15 — Comparison
    comparison_slide(prs, 15)

    # 16 — Demo dashboard (Platform)
    full_image_slide(prs, PLAT / "slide_10.png")

    # 17 — Demo logistics/containers (Platform)
    full_image_slide(prs, PLAT / "slide_11.png")

    # 18 — Story walkthrough (Platform)
    full_image_slide(prs, PLAT / "slide_12.png")

    # 19 — Results & business value
    cards_row(
        prs,
        "Results & Business Value",
        "What Predictive Export Intelligence delivers",
        [
            (
                "Results Observed",
                [
                    "Linked price, demand, route, and container decisions",
                    "Practical limited-capacity container plans",
                    "Management-ready dashboard view",
                    "Explainable via Agent Reasoning + RAG",
                ],
                TEAL,
            ),
            (
                "Business Value",
                [
                    "Faster export decision cycle",
                    "Lower manual market-analysis effort",
                    "Profit visibility before shipment commit",
                    "Clear leadership communication",
                ],
                GOLD,
            ),
            (
                "Demo Outcomes",
                [
                    "One-click four-phase pipeline",
                    "Price → demand → logistics → containers",
                    "Corridors ranked by INR net profit",
                    "Export-first container recommendation",
                ],
                CORAL,
            ),
        ],
        19,
        photo=CONT,
    )

    # 20 — Future roadmap (Platform visual)
    full_image_slide(prs, PLAT / "slide_13.png")

    # 21 — Thank you
    thank_you_slide(prs, 21)

    out = OUT
    try:
        prs.save(str(out))
    except PermissionError:
        out = ROOT / "ExportIntel_AI_Final_v2.pptx"
        prs.save(str(out))

    try:
        prs.save(str(OUT_DL))
        print(f"Also saved: {OUT_DL}")
    except Exception as e:
        print(f"Downloads copy skipped: {e}")

    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
