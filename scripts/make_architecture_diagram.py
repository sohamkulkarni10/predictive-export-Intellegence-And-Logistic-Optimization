"""ExportIntel AI — presentation architecture diagram (16:9)."""

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
OUT_PNG = ROOT / "assets" / "ppt" / "architecture_explain.png"
OUT_PPT = ROOT / "ExportIntel_AI_Architecture.pptx"
OUT_DL_PNG = Path(r"c:\Users\Lenovo\Downloads\ExportIntel_AI_Architecture.png")
OUT_DL_PPT = Path(r"c:\Users\Lenovo\Downloads\ExportIntel_AI_Architecture.pptx")

NAVY = "#081228"
NAVY2 = "#0c1c3a"
PANEL = "#102444"
TEAL = "#1ebec8"
GOLD = "#e6aa37"
CORAL = "#e65a46"
WHITE = "#f4f7fb"
MUTED = "#9eb0c4"
LINE = "#3a5a7a"


def box(ax, x, y, w, h, fc, ec, lw=1.4, r=0.08):
    p = FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle=f"round,pad=0.012,rounding_size={r}",
        facecolor=fc,
        edgecolor=ec,
        linewidth=lw,
        mutation_aspect=0.6,
    )
    ax.add_patch(p)
    return p


def arrow(ax, x1, y1, x2, y2, color=TEAL, lw=1.8):
    ax.add_patch(
        FancyArrowPatch(
            (x1, y1),
            (x2, y2),
            arrowstyle="-|>",
            mutation_scale=12,
            linewidth=lw,
            color=color,
            shrinkA=0,
            shrinkB=0,
        )
    )


def text(ax, x, y, s, size=10, color=WHITE, weight="medium", ha="center", va="center"):
    ax.text(x, y, s, fontsize=size, color=color, ha=ha, va=va, fontweight=weight, fontfamily="DejaVu Sans")


def main():
    fig, ax = plt.subplots(figsize=(19.2, 10.8), dpi=140)
    fig.patch.set_facecolor(NAVY)
    ax.set_facecolor(NAVY)
    ax.set_xlim(0, 19.2)
    ax.set_ylim(0, 10.8)
    ax.axis("off")

    text(ax, 9.6, 10.42, "ExportIntel AI  ·  System Architecture", 22, WHITE, "bold")
    text(
        ax,
        9.6,
        10.02,
        "Custom sequential pipeline  ·  NOT LangChain  ·  React  →  Flask  →  Models / DBs",
        12,
        TEAL,
        "medium",
    )

    # ===== LAYER 1 FRONTEND =====
    box(ax, 0.35, 8.15, 18.5, 1.65, "#0e2248", TEAL, 1.8, 0.1)
    text(ax, 1.55, 9.45, "LAYER 1", 8, GOLD, "bold", ha="left")
    text(ax, 1.55, 9.12, "Frontend  ·  React + Vite  (:5173)", 13, WHITE, "bold", ha="left")
    text(ax, 1.55, 8.78, "User only clicks / views. No ML runs in browser.", 9.5, MUTED, ha="left")

    ui = [
        ("Login", CORAL),
        ("Dashboard", TEAL),
        ("Demand", TEAL),
        ("Price", GOLD),
        ("Logistics", CORAL),
        ("Containers", TEAL),
        ("AI Assistant", GOLD),
        ("Agent Reasoning", MUTED),
    ]
    for i, (lab, col) in enumerate(ui):
        x = 0.6 + i * 2.25
        box(ax, x, 8.32, 2.1, 0.42, PANEL, col, 1.2, 0.06)
        text(ax, x + 1.05, 8.53, lab, 9, WHITE, "bold")

    # arrows UI <-> API
    arrow(ax, 5.2, 8.15, 5.2, 7.55, GOLD)
    text(ax, 4.15, 7.82, "HTTP JSON\nPOST /api/pipeline", 8, GOLD, ha="right")
    arrow(ax, 14.0, 7.55, 14.0, 8.15, TEAL)
    text(ax, 15.15, 7.82, "Predictions\nto UI", 8, TEAL, ha="left")

    # ===== LAYER 2 FLASK =====
    box(ax, 0.35, 4.55, 18.5, 2.9, "#0e2248", GOLD, 1.8, 0.1)
    text(ax, 1.55, 7.15, "LAYER 2", 8, GOLD, "bold", ha="left")
    text(ax, 1.55, 6.82, "Backend  ·  Flask REST API  (:5001)", 13, WHITE, "bold", ha="left")
    text(
        ax,
        1.55,
        6.52,
        "One click → run_pipeline() calls each Python stage in order. Fixed flow, no tool-choosing LLM.",
        9.5,
        MUTED,
        ha="left",
    )

    # auth + rag side cards
    box(ax, 0.55, 4.75, 2.35, 1.55, PANEL, MUTED, 1.2, 0.07)
    text(ax, 1.72, 5.95, "Auth", 11, WHITE, "bold")
    text(ax, 1.72, 5.55, "POST /api/login", 8.5, MUTED)
    text(ax, 1.72, 5.25, "Token check", 8.5, MUTED)
    text(ax, 1.72, 4.98, "before pipeline", 8.5, MUTED)

    box(ax, 16.3, 4.75, 2.35, 1.55, PANEL, GOLD, 1.2, 0.07)
    text(ax, 17.47, 5.95, "RAG path", 11, GOLD, "bold")
    text(ax, 17.47, 5.55, "POST /api/rag/ask", 8.5, MUTED)
    text(ax, 17.47, 5.25, "separate from", 8.5, MUTED)
    text(ax, 17.47, 4.98, "Run Analysis", 8.5, MUTED)

    # pipeline strip
    stages = [
        ("1. News", "GDELT + Google RSS"),
        ("2. Demand", "joblib score 0–1"),
        ("3. Price", "XGBoost INR/qtl"),
        ("4. Logistics", "sell−buy−cost"),
        ("5. Containers", "allocate 6×20FT"),
        ("6. Explain", "Groq reasoning"),
    ]
    for i, (title, sub) in enumerate(stages):
        x = 3.1 + i * 2.15
        col = [TEAL, TEAL, GOLD, CORAL, TEAL, MUTED][i]
        box(ax, x, 4.85, 2.0, 1.35, PANEL, col, 1.5, 0.07)
        text(ax, x + 1.0, 5.75, title, 10.5, col, "bold")
        text(ax, x + 1.0, 5.28, sub, 8.2, WHITE)
        if i < 5:
            arrow(ax, x + 2.02, 5.52, x + 2.13, 5.52, LINE)

    text(ax, 9.6, 4.72, "Custom multi-agent pipeline in pipeline.py   ·   NOT LangChain", 9, GOLD, "bold")

    # ===== LAYER 3 DATA =====
    box(ax, 0.35, 0.28, 18.5, 4.05, "#0e2248", CORAL, 1.8, 0.1)
    text(ax, 1.55, 4.02, "LAYER 3", 8, GOLD, "bold", ha="left")
    text(ax, 1.55, 3.68, "Models · Datasets · Storage", 13, WHITE, "bold", ha="left")

    stores = [
        ("Demand Model", "demand_model_bundle.joblib\n+ Groq news features\n→ opportunity scores", TEAL),
        ("Price Model", "XGBoost generalized.pkl\n+ monthly_price.csv\n→ next-month INR/qtl", GOLD),
        ("Logistics CSVs", "ports · freight · trade\nsell prices · costs\n→ profitable routes", CORAL),
        ("SQLite", "export_ai.db\npipeline_runs\nlogistics_costs", TEAL),
        ("Chroma DB", "trusted CSV chunks\n+ GDELT headlines\n→ RAG retrieval", GOLD),
        ("Groq LLM", "Llama 3.3 70B\nExplain Agent\n+ RAG answers", MUTED),
    ]
    for i, (title, body, col) in enumerate(stores):
        x = 0.55 + i * 3.05
        box(ax, x, 0.48, 2.9, 2.95, PANEL, col, 1.4, 0.08)
        text(ax, x + 1.45, 3.05, title, 11, col, "bold")
        text(ax, x + 1.45, 1.85, body, 9, WHITE)

    # small connectors from pipeline to stores
    arrow(ax, 5.15, 4.85, 2.0, 3.45, TEAL)      # demand stage -> demand model
    arrow(ax, 7.3, 4.85, 5.05, 3.45, GOLD)      # price -> xgb
    arrow(ax, 9.45, 4.85, 8.1, 3.45, CORAL)     # logistics -> csv
    arrow(ax, 11.6, 4.85, 11.15, 3.45, TEAL)    # containers/pipeline -> sqlite
    arrow(ax, 17.47, 4.75, 17.47, 3.45, GOLD)   # rag path -> chroma/groq area

    OUT_PNG.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(OUT_PNG, facecolor=fig.get_facecolor(), bbox_inches="tight", pad_inches=0.12)
    try:
        fig.savefig(OUT_DL_PNG, facecolor=fig.get_facecolor(), bbox_inches="tight", pad_inches=0.12)
    except Exception as e:
        print("Downloads PNG skipped:", e)
    plt.close(fig)

    # one-slide PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(str(OUT_PNG), 0, 0, prs.slide_width, prs.slide_height)
    for path in (OUT_PPT, OUT_DL_PPT):
        try:
            prs.save(str(path))
            print("Saved", path)
        except PermissionError:
            alt = path.with_name(path.stem + "_new" + path.suffix)
            prs.save(str(alt))
            print("Locked — saved", alt)
    print("PNG", OUT_PNG)


if __name__ == "__main__":
    main()
